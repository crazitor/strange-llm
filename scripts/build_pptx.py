#!/usr/bin/env python3
"""
build_pptx.py — 将 slides_spec_v2 的 markdown 规格文件转化为 .pptx
用法: python3 build_pptx.py
"""

import re
import os
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# ─── 路径配置 ───────────────────────────────────
PROJECT = Path("/Users/mik/strange LLM")
OUTPUT = PROJECT / "output"
SCREENSHOTS = OUTPUT / "screenshots"
SPEC_FILES = [
    OUTPUT / "slides_spec_v2_part1.md",
    OUTPUT / "slides_spec_v2_part2.md",
    OUTPUT / "slides_spec_v2_part3.md",
]
OUTPUT_PPTX = OUTPUT / "presentation_v2.pptx"

# ─── 设计参数 ───────────────────────────────────
SLIDE_WIDTH = Inches(13.333)   # 16:9 宽屏
SLIDE_HEIGHT = Inches(7.5)
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # 深蓝黑
BG_SECTION = RGBColor(0x0D, 0x0D, 0x1A)    # 更深的幕标题背景
ACCENT = RGBColor(0x4E, 0xC9, 0xB0)         # 青绿色强调
ACCENT2 = RGBColor(0xFF, 0xA5, 0x00)        # 橙色强调
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
TEXT_DIM = RGBColor(0x88, 0x88, 0x99)
QUOTE_BG = RGBColor(0x14, 0x14, 0x28)       # 金句卡背景

FONT_TITLE = "PingFang SC"
FONT_BODY = "PingFang SC"
FONT_FALLBACK = "Microsoft YaHei"


# ─── 数据结构 ───────────────────────────────────
@dataclass
class SlideSpec:
    page_num: int = 0
    title: str = ""
    section: str = ""
    duration: str = ""
    cumulative: str = ""
    ppt_title: str = ""
    ppt_subtitle_points: list = field(default_factory=list)
    image_type: str = ""        # screenshot / quote-card / text-only / diagram-placeholder
    image_value: str = ""       # filename or description
    speech: str = ""
    speaker_note: str = ""


# ─── Markdown 解析 ──────────────────────────────
def parse_spec_file(filepath: Path) -> list[SlideSpec]:
    """解析一个 slides_spec markdown 文件，返回 SlideSpec 列表"""
    text = filepath.read_text(encoding="utf-8")
    slides = []

    # 按 "### 页 N:" 分割
    page_pattern = re.compile(r'^### 页 (\d+): (.+)$', re.MULTILINE)
    matches = list(page_pattern.finditer(text))

    for i, match in enumerate(matches):
        start = match.start()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        block = text[start:end]

        spec = SlideSpec()
        spec.page_num = int(match.group(1))
        spec.title = match.group(2).strip()

        # 章节
        m = re.search(r'\*\*章节\*\*:\s*(.+)', block)
        if m:
            spec.section = m.group(1).strip()

        # 时间
        m = re.search(r'\*\*时间\*\*:\s*(.+)', block)
        if m:
            spec.duration = m.group(1).strip()

        # 累计时间
        m = re.search(r'\*\*累计时间\*\*:\s*(.+)', block)
        if m:
            spec.cumulative = m.group(1).strip()

        # PPT显示内容
        ppt_section = re.search(
            r'\*\*PPT显示内容\*\*:\s*\n(.*?)(?=\n\*\*演讲词\*\*)',
            block, re.DOTALL
        )
        if ppt_section:
            ppt_block = ppt_section.group(1)

            # 标题
            m = re.search(r'- 标题:\s*(.+)', ppt_block)
            if m:
                spec.ppt_title = m.group(1).strip()

            # 副标题/要点
            points = []
            in_points = False
            for line in ppt_block.split('\n'):
                if '副标题/要点' in line or '副标题' in line:
                    in_points = True
                    # 检查同一行是否有内容
                    after = line.split(':', 1)
                    if len(after) > 1 and after[1].strip():
                        points.append(after[1].strip())
                    continue
                if in_points and line.strip().startswith('-'):
                    point = line.strip().lstrip('-').strip()
                    points.append(point)
                elif in_points and line.strip().startswith('·'):
                    point = line.strip().lstrip('·').strip()
                    points.append(point)
                elif in_points and not line.strip():
                    continue
                elif in_points and '配图' in line:
                    in_points = False
            spec.ppt_subtitle_points = points

            # 配图
            m = re.search(r'- 配图:\s*(.+)', ppt_block)
            if not m:
                # 有时配图在下一行
                m = re.search(r'配图:\s*(.+)', ppt_block)
            if m:
                img_val = m.group(1).strip()
                if img_val.startswith('video') and img_val.endswith('.jpg'):
                    spec.image_type = "screenshot"
                    spec.image_value = img_val
                elif 'quote-card' in img_val.lower() or 'quote_card' in img_val.lower():
                    spec.image_type = "quote-card"
                    # 提取描述
                    desc = re.search(r'[—–-]\s*(.+?)[\]】]', img_val)
                    spec.image_value = desc.group(1) if desc else img_val
                elif 'diagram-placeholder' in img_val.lower() or 'diagram' in img_val.lower():
                    spec.image_type = "diagram-placeholder"
                    desc = re.search(r'[—–-]\s*(.+?)[\]】]', img_val)
                    spec.image_value = desc.group(1) if desc else img_val
                elif 'text-only' in img_val.lower() or 'text_only' in img_val.lower():
                    spec.image_type = "text-only"
                    desc = re.search(r'[—–-]\s*(.+?)[\]】]', img_val)
                    spec.image_value = desc.group(1) if desc else img_val
                else:
                    # 可能是截图文件名但格式不同
                    if '.jpg' in img_val or '.png' in img_val:
                        spec.image_type = "screenshot"
                        # 提取文件名
                        fname = re.search(r'(video\d+_t\d+_.+?\.(jpg|png))', img_val)
                        spec.image_value = fname.group(1) if fname else img_val
                    else:
                        spec.image_type = "text-only"
                        spec.image_value = img_val

        # 演讲词
        speech_match = re.search(
            r'\*\*演讲词\*\*:\s*\n(.*?)(?=\n\*\*讲者提醒\*\*)',
            block, re.DOTALL
        )
        if speech_match:
            speech_lines = speech_match.group(1).strip()
            # 去掉 > 引用符号
            speech_lines = re.sub(r'^>\s?', '', speech_lines, flags=re.MULTILINE)
            spec.speech = speech_lines.strip()

        # 讲者提醒
        m = re.search(r'\*\*讲者提醒\*\*:\s*(.+?)(?:\n---|\n\n|\Z)', block, re.DOTALL)
        if m:
            spec.speaker_note = m.group(1).strip()

        slides.append(spec)

    return slides


# ─── PPTX 构建 ──────────────────────────────────
def set_slide_bg(slide, color: RGBColor):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_BODY):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_points(slide, left, top, width, height, points,
                      font_size=16, color=TEXT_LIGHT):
    """添加要点列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # 清理要点文本
        clean_point = point.strip()
        if clean_point.startswith(('1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.')):
            p.text = clean_point
        else:
            p.text = f"• {clean_point}"

        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(6)

    return txBox


def build_section_slide(prs, spec: SlideSpec):
    """构建幕标题页（深色背景大标题）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_bg(slide, BG_SECTION)

    # 大标题
    add_text_box(
        slide, Inches(1), Inches(2), Inches(11), Inches(2),
        spec.ppt_title or spec.title,
        font_size=44, color=ACCENT, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE
    )

    # 副标题
    if spec.ppt_subtitle_points:
        subtitle_text = "\n".join(spec.ppt_subtitle_points)
        add_text_box(
            slide, Inches(2), Inches(4.2), Inches(9), Inches(1.5),
            subtitle_text,
            font_size=20, color=TEXT_LIGHT,
            alignment=PP_ALIGN.CENTER
        )

    # 章节标签
    add_text_box(
        slide, Inches(1), Inches(6.5), Inches(4), Inches(0.5),
        spec.section,
        font_size=12, color=TEXT_DIM
    )

    # 时间标签
    add_text_box(
        slide, Inches(9), Inches(6.5), Inches(3), Inches(0.5),
        f"累计 {spec.cumulative}",
        font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.RIGHT
    )

    # 演讲词放入备注
    notes = slide.notes_slide
    notes_tf = notes.notes_text_frame
    notes_tf.text = f"【演讲词】\n{spec.speech}\n\n【讲者提醒】\n{spec.speaker_note}"

    return slide


def build_quote_card(prs, spec: SlideSpec):
    """构建金句卡片页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, QUOTE_BG)

    # 提取金句（通常在演讲词或image_value中）
    quote_text = spec.image_value if spec.image_value else spec.ppt_title

    # 大字金句
    add_text_box(
        slide, Inches(1.5), Inches(1.5), Inches(10), Inches(4),
        ("\u201C" + quote_text + "\u201D") if not quote_text.startswith("\u201C") else quote_text,
        font_size=36, color=TEXT_WHITE, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE
    )

    # 如果有要点，在下方小字展示
    if spec.ppt_subtitle_points:
        add_bullet_points(
            slide, Inches(2), Inches(5), Inches(9), Inches(1.5),
            spec.ppt_subtitle_points[:3],
            font_size=14, color=TEXT_DIM
        )

    # 章节 + 时间
    add_text_box(
        slide, Inches(1), Inches(6.5), Inches(4), Inches(0.5),
        spec.section, font_size=12, color=TEXT_DIM
    )
    add_text_box(
        slide, Inches(9), Inches(6.5), Inches(3), Inches(0.5),
        f"累计 {spec.cumulative}",
        font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.RIGHT
    )

    # 备注
    notes = slide.notes_slide
    notes_tf = notes.notes_text_frame
    notes_tf.text = f"【演讲词】\n{spec.speech}\n\n【讲者提醒】\n{spec.speaker_note}"

    return slide


def build_content_slide(prs, spec: SlideSpec, has_image=False, img_path=None):
    """构建内容页（标题 + 要点 + 可选图片）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    if has_image and img_path and img_path.exists():
        # 左侧文字，右侧图片布局
        # 标题
        add_text_box(
            slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.8),
            spec.ppt_title or spec.title,
            font_size=28, color=ACCENT, bold=True, font_name=FONT_TITLE
        )

        # 要点
        if spec.ppt_subtitle_points:
            add_bullet_points(
                slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(4.5),
                spec.ppt_subtitle_points,
                font_size=15, color=TEXT_LIGHT
            )

        # 图片（右侧）
        try:
            img = Image.open(img_path)
            img_w, img_h = img.size
            aspect = img_w / img_h

            max_w = Inches(6)
            max_h = Inches(5.5)

            if aspect > max_w / max_h:
                w = max_w
                h = int(w / aspect)
            else:
                h = max_h
                w = int(h * aspect)

            left = Inches(13.333) - w - Inches(0.4)
            top = Inches(1.2)
            slide.shapes.add_picture(str(img_path), left, top, w, h)
        except Exception as e:
            # 图片加载失败，添加占位文本
            add_text_box(
                slide, Inches(7), Inches(2.5), Inches(5.5), Inches(2),
                f"[图片: {img_path.name}]",
                font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER
            )
    else:
        # 全宽文字布局
        # 标题
        add_text_box(
            slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
            spec.ppt_title or spec.title,
            font_size=32, color=ACCENT, bold=True, font_name=FONT_TITLE
        )

        # 要点
        if spec.ppt_subtitle_points:
            add_bullet_points(
                slide, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                spec.ppt_subtitle_points,
                font_size=18, color=TEXT_LIGHT
            )

        # 如果是 diagram-placeholder，添加占位区域
        if spec.image_type == "diagram-placeholder":
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(2), Inches(4), Inches(9), Inches(2.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)
            shape.line.color.rgb = ACCENT
            shape.line.width = Pt(1)

            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"[待设计] {spec.image_value}"
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_DIM
            p.alignment = PP_ALIGN.CENTER

    # 章节 + 累计时间底栏
    add_text_box(
        slide, Inches(0.5), Inches(6.8), Inches(6), Inches(0.4),
        f"{spec.section}  |  {spec.duration}",
        font_size=10, color=TEXT_DIM
    )
    add_text_box(
        slide, Inches(9), Inches(6.8), Inches(3.5), Inches(0.4),
        f"页 {spec.page_num}  |  累计 {spec.cumulative}",
        font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.RIGHT
    )

    # 备注
    notes = slide.notes_slide
    notes_tf = notes.notes_text_frame
    notes_tf.text = f"【演讲词】\n{spec.speech}\n\n【讲者提醒】\n{spec.speaker_note}"

    return slide


def is_section_title(spec: SlideSpec) -> bool:
    """判断是否为幕标题页"""
    title_lower = spec.title.lower()
    return (
        '标题页' in title_lower or
        '第一幕' in spec.title or '第二幕' in spec.title or
        '第三幕' in spec.title or '第四幕' in spec.title or
        '第五幕' in spec.title or '第六幕' in spec.title or
        '第七幕' in spec.title or '序幕' in spec.title
    ) and spec.image_type in ('text-only', '')


def build_presentation(all_slides: list[SlideSpec]) -> Presentation:
    """构建完整的演示文稿"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    total = len(all_slides)
    for i, spec in enumerate(all_slides):
        print(f"  [{i+1}/{total}] 页 {spec.page_num}: {spec.title[:30]}... ({spec.image_type})")

        # 根据类型选择构建方式
        if is_section_title(spec):
            build_section_slide(prs, spec)
        elif spec.image_type == "quote-card":
            build_quote_card(prs, spec)
        elif spec.image_type == "screenshot":
            img_path = SCREENSHOTS / spec.image_value
            build_content_slide(prs, spec, has_image=True, img_path=img_path)
        elif spec.image_type == "diagram-placeholder":
            build_content_slide(prs, spec, has_image=False)
        else:
            # text-only 或其他
            build_content_slide(prs, spec, has_image=False)

    return prs


# ─── 主流程 ─────────────────────────────────────
def main():
    print("=" * 60)
    print("PPTX Build — AI 科普演讲 v2")
    print("=" * 60)

    # 1. 解析所有 spec 文件
    all_slides = []
    for spec_file in SPEC_FILES:
        print(f"\n解析: {spec_file.name}")
        slides = parse_spec_file(spec_file)
        print(f"  → {len(slides)} 页")
        all_slides.extend(slides)

    print(f"\n总计: {len(all_slides)} 页")

    # 2. 统计配图类型
    type_counts = {}
    for s in all_slides:
        t = s.image_type or "unknown"
        type_counts[t] = type_counts.get(t, 0) + 1
    print("\n配图类型统计:")
    for t, c in sorted(type_counts.items()):
        print(f"  {t}: {c}")

    # 3. 检查截图文件存在性
    missing = []
    for s in all_slides:
        if s.image_type == "screenshot":
            path = SCREENSHOTS / s.image_value
            if not path.exists():
                missing.append((s.page_num, s.image_value))
    if missing:
        print(f"\n⚠️ {len(missing)} 张截图文件缺失:")
        for pn, fn in missing[:10]:
            print(f"  页 {pn}: {fn}")

    # 4. 构建 PPTX
    print("\n开始构建 PPTX...")
    prs = build_presentation(all_slides)

    # 5. 保存
    prs.save(str(OUTPUT_PPTX))
    size_mb = OUTPUT_PPTX.stat().st_size / (1024 * 1024)
    print(f"\n✅ 已保存: {OUTPUT_PPTX}")
    print(f"   文件大小: {size_mb:.1f} MB")
    print(f"   总页数: {len(all_slides)}")
    print("=" * 60)


if __name__ == "__main__":
    main()
