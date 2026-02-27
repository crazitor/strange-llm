#!/usr/bin/env python3
"""
validate_pptx_v3.py — PPT v3 自动验收脚本

检查项：
  1. Spec页数 = PPTX页数
  2. 每页notes非空
  3. visual引用的文件存在
  4. on-screen内容行数/字符数不超阈值
  5. 禁止残留标记
  6. 字号层级覆盖 >= 4 档
  7. 页面类型分布报告

用法：python3.11 scripts/validate_pptx_v3.py
"""

import re
import sys
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional

# ─── 路径配置 ───────────────────────────────────────
PROJECT = Path("/Users/mik/strange LLM")
OUTPUT = PROJECT / "output"
SCREENSHOTS = OUTPUT / "screenshots"
DIAGRAMS = PROJECT / "assets" / "diagrams" / "rendered"
SPEC_FILES = [
    OUTPUT / "slides_spec_v3_part1.md",
    OUTPUT / "slides_spec_v3_part2.md",
]
PPTX_FILE = OUTPUT / "presentation_v3.pptx"

# ─── ANSI 颜色 ─────────────────────────────────────
GREEN = "\033[92m"
RED = "\033[91m"
YELLOW = "\033[93m"
CYAN = "\033[96m"
BOLD = "\033[1m"
DIM = "\033[2m"
RESET = "\033[0m"


def ok(msg: str) -> str:
    return f"{GREEN}✓{RESET} {msg}"


def fail(msg: str) -> str:
    return f"{RED}✗{RESET} {msg}"


def warn(msg: str) -> str:
    return f"{YELLOW}⚠{RESET} {msg}"


def header(msg: str) -> str:
    return f"\n{BOLD}{CYAN}─── {msg} ───{RESET}"


# ─── Spec 数据结构 ─────────────────────────────────
@dataclass
class SlideSpec:
    page_num: int = 0
    title: str = ""
    slide_type: str = ""
    visual: str = ""
    on_screen: str = ""
    notes: str = ""
    raw_block: str = ""


# ─── Spec 解析 ─────────────────────────────────────
def parse_spec_files(spec_files: list[Path]) -> list[SlideSpec]:
    """解析所有spec文件，返回SlideSpec列表。"""
    slides: list[SlideSpec] = []

    for filepath in spec_files:
        if not filepath.exists():
            print(warn(f"Spec文件不存在: {filepath.name}"))
            continue

        text = filepath.read_text(encoding="utf-8")
        # 匹配页面标题行: ### 页 N: ...  或  ### P N: ...
        page_pattern = re.compile(r'^###\s+(?:页|P)\s*(\d+)\s*[:：]\s*(.+)$', re.MULTILINE)
        matches = list(page_pattern.finditer(text))

        for i, match in enumerate(matches):
            start = match.start()
            end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
            block = text[start:end]

            spec = SlideSpec()
            spec.page_num = int(match.group(1))
            spec.title = match.group(2).strip()
            spec.raw_block = block

            # type 字段: **type**: xxx 或 **类型**: xxx
            m = re.search(r'\*\*type\*\*\s*[:：]\s*(.+)', block, re.IGNORECASE)
            if not m:
                m = re.search(r'\*\*类型\*\*\s*[:：]\s*(.+)', block)
            if m:
                spec.slide_type = m.group(1).strip().strip('`').lower()

            # visual 字段
            m = re.search(r'\*\*visual\*\*\s*[:：]\s*(.+)', block, re.IGNORECASE)
            if not m:
                m = re.search(r'\*\*配图\*\*\s*[:：]\s*(.+)', block)
            if not m:
                m = re.search(r'- 配图\s*[:：]\s*(.+)', block)
            if m:
                spec.visual = m.group(1).strip()

            # on-screen 字段 (通常单行: - **on-screen**: 内容)
            m = re.search(
                r'\*\*on-screen\*\*\s*[:：]\s*(.+)',
                block, re.IGNORECASE
            )
            if not m:
                # 尝试提取 PPT显示内容 区域 (多行)
                m = re.search(
                    r'\*\*PPT显示内容\*\*\s*[:：]?\s*\n(.*?)(?=\n-?\s*\*\*|\n###|\Z)',
                    block, re.DOTALL | re.IGNORECASE
                )
            if m:
                spec.on_screen = m.group(1).strip()

            # notes 字段 (演讲词/speaker notes)
            m = re.search(
                r'\*\*(?:notes|演讲词|讲者提醒)\*\*\s*[:：]?\s*\n?(.*?)(?=\n\*\*|\n###|\n---|\Z)',
                block, re.DOTALL | re.IGNORECASE
            )
            if m:
                spec.notes = m.group(1).strip()

            slides.append(spec)

    # 按页码排序
    slides.sort(key=lambda s: s.page_num)
    return slides


def extract_visual_files(slides: list[SlideSpec]) -> list[tuple[int, str, str]]:
    """从visual字段提取引用的文件名。
    返回 [(page_num, filename, raw_visual), ...]
    """
    results = []
    for s in slides:
        if not s.visual:
            continue
        raw = s.visual

        # 跳过纯文字描述类visual（无文件引用）
        # 匹配各种图片文件名
        filenames = re.findall(r'([\w\-]+\.(?:png|jpg|jpeg|svg|webp))', raw, re.IGNORECASE)
        for fname in filenames:
            results.append((s.page_num, fname, raw))

    return results


# ─── On-screen 内容检查 ────────────────────────────
# 阈值定义
ONSCREEN_LIMITS = {
    "key-phrase":    {"max_lines": 1, "max_chars": 40},
    "quote-impact":  {"max_lines": 1, "max_chars": 35},
    "full-visual":   {"max_lines": 2, "max_chars": 35},
    "data-callout":  {"max_lines": 2, "max_chars_num_line": 10},
    "evidence":      {"max_lines": 4},
    "split-compare": {"max_labels": 4},
    "bridge":        {"max_sentences": 2},
}


def count_effective_lines(text: str) -> tuple[int, int]:
    """返回(行数, 最长行字符数)，忽略空行和markdown标记行。"""
    lines = []
    for line in text.split("\n"):
        stripped = line.strip()
        # 跳过空行、纯markdown标记行
        if not stripped:
            continue
        # 去除 markdown 前缀（- / * / 数字. 等）
        cleaned = re.sub(r'^[\-\*\d\.]+\s*', '', stripped)
        if cleaned:
            lines.append(cleaned)

    if not lines:
        return 0, 0
    max_chars = max(len(l) for l in lines)
    return len(lines), max_chars


def check_onscreen_limits(slides: list[SlideSpec]) -> list[str]:
    """检查on-screen内容是否超过阈值。返回问题列表。"""
    issues = []

    for s in slides:
        if not s.on_screen or not s.slide_type:
            continue

        stype = s.slide_type
        limits = ONSCREEN_LIMITS.get(stype)
        if not limits:
            continue

        nlines, max_chars = count_effective_lines(s.on_screen)

        # 检查行数
        if "max_lines" in limits and nlines > limits["max_lines"]:
            excerpt = s.on_screen[:60].replace("\n", "\\n")
            issues.append(
                f"P{s.page_num} {stype}: {nlines}行 > {limits['max_lines']}行上限 "
                f'"{excerpt}..."'
            )
            continue

        # 检查字符数
        if "max_chars" in limits and max_chars > limits["max_chars"]:
            excerpt = s.on_screen[:60].replace("\n", "\\n")
            issues.append(
                f"P{s.page_num} {stype}: 最长行{max_chars}字 > {limits['max_chars']}字上限 "
                f'"{excerpt}..."'
            )
            continue

        # data-callout 数字行检查
        if "max_chars_num_line" in limits:
            for line in s.on_screen.split("\n"):
                stripped = line.strip().lstrip("-*· ")
                # 检测数字行（主要由数字/符号/百分号组成）
                if stripped and re.match(r'^[\d\.\,\%\+\-×\/\s→≈≤≥<>]+$', stripped):
                    if len(stripped) > limits["max_chars_num_line"]:
                        issues.append(
                            f"P{s.page_num} {stype}: 数字行{len(stripped)}字 > "
                            f"{limits['max_chars_num_line']}字上限 \"{stripped}\""
                        )

        # split-compare 标签数检查
        if "max_labels" in limits:
            label_lines = [
                l.strip() for l in s.on_screen.split("\n")
                if l.strip() and not l.strip().startswith("#")
            ]
            if len(label_lines) > limits["max_labels"]:
                issues.append(
                    f"P{s.page_num} {stype}: {len(label_lines)}个标签 > "
                    f"{limits['max_labels']}个上限"
                )

        # bridge 句数检查
        if "max_sentences" in limits:
            text = s.on_screen.strip()
            sentences = re.split(r'[。！？.!?]', text)
            sentences = [sent for sent in sentences if sent.strip()]
            if len(sentences) > limits["max_sentences"]:
                excerpt = text[:50].replace("\n", "\\n")
                issues.append(
                    f"P{s.page_num} {stype}: {len(sentences)}句 > "
                    f"{limits['max_sentences']}句上限 \"{excerpt}...\""
                )

    return issues


# ─── 残留检查 ──────────────────────────────────────
FORBIDDEN_PATTERNS = [
    (r'\[设计指示：', "[设计指示：...]"),
    (r'配图：', "配图："),
    (r'diagram-placeholder', "diagram-placeholder"),
]


def check_forbidden(slides: list[SlideSpec]) -> list[tuple[str, list[int]]]:
    """检查spec中是否有禁止残留的字符串。返回 [(pattern_desc, [page_nums]), ...]"""
    results = []
    for pattern_re, desc in FORBIDDEN_PATTERNS:
        found_pages = []
        for s in slides:
            if re.search(pattern_re, s.raw_block):
                found_pages.append(s.page_num)
        results.append((desc, found_pages))
    return results


# ─── PPTX 检查 ─────────────────────────────────────
def check_pptx(pptx_path: Path, expected_count: int) -> dict:
    """检查PPTX文件。返回检查结果字典。"""
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx 未安装"}

    if not pptx_path.exists():
        return {"error": f"PPTX文件不存在: {pptx_path.name}"}

    prs = Presentation(str(pptx_path))
    total = len(prs.slides)

    # notes 检查
    notes_ok = 0
    notes_empty_pages = []
    for i, slide in enumerate(prs.slides, 1):
        has_notes = False
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                has_notes = True
        if has_notes:
            notes_ok += 1
        else:
            notes_empty_pages.append(i)

    # 字号收集 — 检查 run 级别和 paragraph 默认级别
    font_sizes = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    # paragraph 级别 (defRPr)
                    if para.font.size is not None:
                        font_sizes.add(para.font.size.pt)
                    for run in para.runs:
                        if run.font.size is not None:
                            font_sizes.add(run.font.size.pt)

    return {
        "total": total,
        "expected": expected_count,
        "match": total == expected_count,
        "notes_ok": notes_ok,
        "notes_empty_pages": notes_empty_pages,
        "font_sizes": sorted(font_sizes),
    }


# ─── 主函数 ────────────────────────────────────────
def main():
    print(f"\n{BOLD}═══ PPT v3 验收报告 ═══{RESET}\n")

    passed = 0
    failed = 0
    fix_items: list[str] = []

    # ── 1. 解析Spec ──
    slides = parse_spec_files(SPEC_FILES)
    spec_count = len(slides)

    if spec_count == 0:
        print(fail(f"Spec页数: 0 (未找到任何页面定义)"))
        print(f"  {DIM}检查文件: {', '.join(f.name for f in SPEC_FILES)}{RESET}")
        failed += 1
    else:
        print(ok(f"Spec页数: {spec_count}"))
        passed += 1

    # ── 2. PPTX检查 ──
    pptx_result = None
    if PPTX_FILE.exists():
        pptx_result = check_pptx(PPTX_FILE, spec_count)

        if "error" in pptx_result:
            print(fail(f"PPTX检查: {pptx_result['error']}"))
            failed += 1
        else:
            pptx_total = pptx_result["total"]
            if pptx_result["match"]:
                print(ok(f"PPTX页数: {pptx_total} (匹配)"))
                passed += 1
            else:
                print(fail(f"PPTX页数: {pptx_total} (期望 {spec_count})"))
                failed += 1
                fix_items.append(f"页数不匹配({pptx_total} vs {spec_count})")

            # Notes覆盖率
            notes_ok = pptx_result["notes_ok"]
            notes_pct = notes_ok / pptx_total * 100 if pptx_total > 0 else 0
            if notes_ok == pptx_total:
                print(ok(f"Notes覆盖: {notes_ok}/{pptx_total} ({notes_pct:.0f}%)"))
                passed += 1
            else:
                n_empty = len(pptx_result["notes_empty_pages"])
                pages_str = ", ".join(str(p) for p in pptx_result["notes_empty_pages"][:10])
                if n_empty > 10:
                    pages_str += f" ...等{n_empty}页"
                print(fail(f"Notes覆盖: {notes_ok}/{pptx_total} ({notes_pct:.0f}%)"))
                print(f"  {DIM}缺notes页: {pages_str}{RESET}")
                failed += 1
                fix_items.append(f"notes缺失({n_empty})")
    else:
        print(warn(f"PPTX文件不存在: {PPTX_FILE.name} (跳过PPTX相关检查)"))

    # ── 3. Visual 文件检查 ──
    print(header("Visual文件检查"))
    visual_refs = extract_visual_files(slides)
    missing_visuals = []

    if not visual_refs:
        print(f"  {DIM}(未在spec中找到visual文件引用){RESET}")
    else:
        for page_num, fname, raw in visual_refs:
            # 在多个位置查找文件
            found = False
            for search_dir in [SCREENSHOTS, DIAGRAMS, OUTPUT]:
                if (search_dir / fname).exists():
                    found = True
                    break

            if found:
                print(f"  {ok(f'{fname}')}")
            else:
                print(f"  {fail(f'{fname} → 缺失 (P{page_num})')}")
                missing_visuals.append((page_num, fname))

    if missing_visuals:
        print(f"  {RED}缺失文件: {len(missing_visuals)}个{RESET}")
        failed += 1
        fix_items.append(f"visual缺失({len(missing_visuals)})")
    else:
        if visual_refs:
            print(f"  {GREEN}全部存在{RESET}")
        passed += 1

    # ── 4. On-screen 内容检查 ──
    print(header("On-screen内容检查"))
    onscreen_issues = check_onscreen_limits(slides)

    # 先打印一些正常样本
    ok_shown = 0
    for s in slides:
        if not s.on_screen or not s.slide_type:
            continue
        stype = s.slide_type
        if stype not in ONSCREEN_LIMITS:
            continue
        # 只打印未出问题的页面样本
        is_ok = not any(f"P{s.page_num} " in issue for issue in onscreen_issues)
        if is_ok and ok_shown < 3:
            nlines, max_chars = count_effective_lines(s.on_screen)
            excerpt = s.on_screen.split("\n")[0].strip()[:40]
            limits = ONSCREEN_LIMITS[stype]
            char_limit = limits.get("max_chars", "—")
            q = '"'
            msg = f"P{s.page_num} {stype}: {q}{excerpt}{q} ({max_chars}字, ≤{char_limit})"
            print(f"  {ok(msg)}")
            ok_shown += 1

    for issue in onscreen_issues:
        print(f"  {fail(issue)}")

    if onscreen_issues:
        print(f"  {RED}超标页面: {len(onscreen_issues)}个{RESET}")
        failed += 1
        fix_items.append(f"on-screen超标({len(onscreen_issues)})")
    else:
        if ok_shown > 0:
            print(f"  {GREEN}全部合规{RESET}")
        passed += 1

    # ── 5. 残留检查 ──
    print(header("残留检查"))
    forbidden_results = check_forbidden(slides)
    has_forbidden = False
    for desc, pages in forbidden_results:
        if pages:
            pages_str = ", ".join(f"P{p}" for p in pages[:10])
            print(f"  {fail(f'发现 {desc} (共{len(pages)}处): {pages_str}')}")
            has_forbidden = True
        else:
            print(f"  {ok(f'未发现 {desc}')}")

    if has_forbidden:
        total_forbidden = sum(len(pages) for _, pages in forbidden_results)
        failed += 1
        fix_items.append(f"残留标记({total_forbidden})")
    else:
        passed += 1

    # ── 6. 页面类型分布 ──
    print(header("页面类型分布"))
    type_counts: dict[str, int] = {}
    for s in slides:
        t = s.slide_type if s.slide_type else "(未标注)"
        type_counts[t] = type_counts.get(t, 0) + 1

    if type_counts:
        max_label_len = max(len(t) for t in type_counts)
        for t, count in sorted(type_counts.items(), key=lambda x: -x[1]):
            pct = count / spec_count * 100 if spec_count > 0 else 0
            bar = "█" * int(pct / 3)
            print(f"  {t:<{max_label_len + 2}} {count:>3} ({pct:4.1f}%) {DIM}{bar}{RESET}")
    else:
        print(f"  {DIM}(无类型信息){RESET}")

    # ── 7. 字号层级 ──
    print(header("字号层级"))
    if pptx_result and "font_sizes" in pptx_result:
        font_sizes = pptx_result["font_sizes"]
        n_sizes = len(font_sizes)
        sizes_str = ", ".join(f"{s:.0f}pt" for s in font_sizes)

        if n_sizes >= 4:
            print(f"  {ok(f'使用了 {n_sizes} 档字号 (≥4)')}")
            print(f"  {DIM}{sizes_str}{RESET}")
            passed += 1
        else:
            print(f"  {fail(f'仅使用了 {n_sizes} 档字号 (<4)')}")
            print(f"  {DIM}{sizes_str}{RESET}")
            failed += 1
            fix_items.append(f"字号层级不足({n_sizes})")
    else:
        print(f"  {DIM}(PPTX不可用，跳过){RESET}")

    # ── 总结 ──
    total_checks = passed + failed
    print(f"\n{BOLD}═══ 总结 ═══{RESET}")
    if passed > 0:
        print(f"{GREEN}✓ 通过: {passed}/{total_checks}{RESET}")
    if failed > 0:
        print(f"{RED}✗ 失败: {failed}/{total_checks}{RESET}")
        fix_str = ", ".join(fix_items)
        print(f"  {YELLOW}需修复: {fix_str}{RESET}")
    elif total_checks > 0:
        print(f"\n  {GREEN}{BOLD}🎉 全部通过！{RESET}")

    print()
    sys.exit(1 if failed > 0 else 0)


if __name__ == "__main__":
    main()
