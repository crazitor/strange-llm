#!/usr/bin/env python3
"""
build_pptx_v3.py — PPTX v3 构建脚本
从 slides_spec_v3 的 markdown 规格文件构建专业级演示文稿

v3 vs v2 关键改进：
- 11种显式声明的页面类型（取代启发式猜测）
- 7级字号层次（10-72pt）
- roadmap + bridge 导航系统
- 顶部3pt进度条
- on-screen 只显示 key phrase，不再是全部 bullet
"""

import re
import os
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional
from copy import deepcopy
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree
from PIL import Image, ImageDraw, ImageFilter


# ─── 路径配置 ───────────────────────────────────────────
PROJECT = Path("/Users/mik/strange LLM")
OUTPUT = PROJECT / "output"
SCREENSHOTS = OUTPUT / "screenshots"
DIAGRAMS_RENDERED = PROJECT / "assets" / "diagrams" / "rendered"
SPEC_FILES = [
    OUTPUT / "slides_spec_v3_part1.md",
    OUTPUT / "slides_spec_v3_part2.md",
]
OUTPUT_PPTX = OUTPUT / "presentation_v3.pptx"

FONT_TITLE = "PingFang SC"
FONT_BODY = "PingFang SC"

# ─── 幻灯片尺寸 ─────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)   # 16:9 宽屏
SLIDE_HEIGHT = Inches(7.5)

# ─── 配色方案 ───────────────────────────────────────────
BG_PRIMARY   = RGBColor(0x0A, 0x0F, 0x1E)   # 主背景（更深）
BG_SECONDARY = RGBColor(0x12, 0x15, 0x2A)   # 渐变终点
BG_SECTION   = RGBColor(0x06, 0x0A, 0x18)   # act-title 更深

ACCENT1      = RGBColor(0x00, 0xD4, 0xAA)   # 青绿
ACCENT2      = RGBColor(0xFF, 0x6B, 0x35)   # 橙

TEXT_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_BODY    = RGBColor(0xE8, 0xE8, 0xE8)
TEXT_DIM     = RGBColor(0x6B, 0x72, 0x80)
TEXT_XDIM    = RGBColor(0x3A, 0x40, 0x52)   # 极暗（未来幕）

QUOTE_BG1    = RGBColor(0x0A, 0x0E, 0x1A)   # 金句渐变起点
QUOTE_BG2    = RGBColor(0x1A, 0x10, 0x30)   # 金句渐变终点（紫调）

PROGRESS_BG  = RGBColor(0x1E, 0x25, 0x3A)   # 进度条背景
BORDER_COLOR = RGBColor(0x2A, 0x35, 0x50)   # 微妙边框

# ─── 字号层次 ───────────────────────────────────────────
FONT_SIZES = {
    'act-title':     60,
    'key-phrase':    44,
    'data-number':   72,
    'quote':         36,
    'section-title': 28,
    'caption':       16,
    'footer':        10,
}

# ─── 进度条 ─────────────────────────────────────────────
PROGRESS_BAR_HEIGHT = Pt(3)
PROGRESS_BAR_Y = 0  # 页面顶部

# ─── 幕信息 ─────────────────────────────────────────────
ACTS = {
    'prologue': {'label': '序幕',   'num': '序', 'color': ACCENT1, 'pages': (1, 6)},
    'act1':     {'label': '第一幕', 'num': '壹', 'color': ACCENT1, 'pages': (7, 21)},
    'act2':     {'label': '第二幕', 'num': '贰', 'color': ACCENT2, 'pages': (22, 33)},
    'act3':     {'label': '第三幕', 'num': '叁', 'color': ACCENT1, 'pages': (34, 50)},
    'act4':     {'label': '第四幕', 'num': '肆', 'color': ACCENT2, 'pages': (51, 63)},
    'act5':     {'label': '第五幕', 'num': '伍', 'color': ACCENT1, 'pages': (64, 75)},
    'act6':     {'label': '第六幕', 'num': '陆', 'color': ACCENT2, 'pages': (76, 83)},
    'act7':     {'label': '第七幕', 'num': '柒', 'color': ACCENT1, 'pages': (84, 90)},
}

# 幕的有序列表（用于 roadmap）
ACT_ORDER = ['prologue', 'act1', 'act2', 'act3', 'act4', 'act5', 'act6', 'act7']


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 数据结构
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
@dataclass
class SlideSpecV3:
    page_num: int = 0
    title: str = ""
    slide_type: str = ""       # 11种类型之一
    act: str = ""              # prologue | act1..act7
    on_screen: str = ""        # 屏幕显示文字
    visual: str = "none"       # 图片文件名 | "none"
    notes: str = ""            # 完整演讲词
    duration: float = 0.0


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Spec 解析器
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def parse_v3_spec(filepath: Path) -> list[SlideSpecV3]:
    """解析 v3 格式的 spec markdown 文件"""
    text = filepath.read_text(encoding="utf-8")
    slides = []

    # 匹配 ### P[N]: [标题]
    page_pattern = re.compile(r'^### P(\d+):\s*(.+)$', re.MULTILINE)
    matches = list(page_pattern.finditer(text))

    for i, match in enumerate(matches):
        start = match.start()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        block = text[start:end]

        spec = SlideSpecV3()
        spec.page_num = int(match.group(1))
        spec.title = match.group(2).strip()

        # type
        m = re.search(r'\*\*type\*\*:\s*(.+)', block)
        if m:
            spec.slide_type = m.group(1).strip().lower()

        # act
        m = re.search(r'\*\*act\*\*:\s*(.+)', block)
        if m:
            spec.act = m.group(1).strip().lower()

        # on-screen — 可能多行
        m = re.search(r'\*\*on-screen\*\*:\s*(.+?)(?=\n- \*\*|\Z)', block, re.DOTALL)
        if m:
            raw = m.group(1).strip()
            # 清理 markdown continuation lines
            lines = []
            for line in raw.split('\n'):
                line = line.strip()
                if line.startswith('- **'):
                    break
                lines.append(line)
            spec.on_screen = '\n'.join(lines).strip()

        # visual
        m = re.search(r'\*\*visual\*\*:\s*(.+)', block)
        if m:
            spec.visual = m.group(1).strip()
            if spec.visual.lower() == '"none"' or spec.visual.lower() == 'none':
                spec.visual = "none"

        # notes — 可能多行
        m = re.search(r'\*\*notes\*\*:\s*(.+?)(?=\n- \*\*|\Z)', block, re.DOTALL)
        if m:
            raw = m.group(1).strip()
            lines = []
            for line in raw.split('\n'):
                stripped = line.strip()
                if stripped.startswith('- **'):
                    break
                lines.append(stripped)
            spec.notes = '\n'.join(lines).strip()

        # Duration
        m = re.search(r'\*\*[Dd]uration\*\*:\s*(.+)', block)
        if m:
            try:
                spec.duration = float(re.search(r'[\d.]+', m.group(1)).group())
            except (AttributeError, ValueError):
                spec.duration = 0.0

        slides.append(spec)

    return slides


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 图片处理工具
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def add_rounded_shadow(img_path: Path, radius=20, shadow_offset=8, shadow_blur=15) -> BytesIO:
    """给图片添加圆角和阴影效果，返回处理后的 BytesIO"""
    img = Image.open(img_path).convert("RGBA")

    # 缩放到合理尺寸
    max_w = 900
    if img.width > max_w:
        ratio = max_w / img.width
        img = img.resize((max_w, int(img.height * ratio)), Image.LANCZOS)

    # 创建圆角蒙版
    mask = Image.new("L", img.size, 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle([(0, 0), img.size], radius=radius, fill=255)

    # 应用圆角
    rounded = Image.new("RGBA", img.size, (0, 0, 0, 0))
    rounded.paste(img, mask=mask)

    # 创建阴影
    canvas_w = img.width + shadow_offset * 2 + shadow_blur * 2
    canvas_h = img.height + shadow_offset * 2 + shadow_blur * 2
    canvas = Image.new("RGBA", (canvas_w, canvas_h), (0, 0, 0, 0))

    shadow = Image.new("RGBA", img.size, (0, 0, 0, 80))
    shadow_masked = Image.new("RGBA", img.size, (0, 0, 0, 0))
    shadow_masked.paste(shadow, mask=mask)

    shadow_canvas = Image.new("RGBA", (canvas_w, canvas_h), (0, 0, 0, 0))
    shadow_canvas.paste(shadow_masked,
                        (shadow_blur + shadow_offset, shadow_blur + shadow_offset))
    shadow_canvas = shadow_canvas.filter(ImageFilter.GaussianBlur(shadow_blur))

    canvas = Image.alpha_composite(canvas, shadow_canvas)
    canvas.paste(rounded, (shadow_blur, shadow_blur), rounded)

    # 细边框
    border_draw = ImageDraw.Draw(canvas)
    border_draw.rounded_rectangle(
        [(shadow_blur, shadow_blur),
         (shadow_blur + img.width - 1, shadow_blur + img.height - 1)],
        radius=radius,
        outline=(0, 212, 170, 60),  # ACCENT1 半透明
        width=2
    )

    buf = BytesIO()
    canvas.save(buf, format="PNG")
    buf.seek(0)
    return buf


def resolve_visual(spec: SlideSpecV3) -> Optional[Path]:
    """解析 visual 字段，找到实际图片路径"""
    if not spec.visual or spec.visual == "none":
        return None

    visual = spec.visual.strip().strip('"').strip("'")

    # 尝试截图目录
    path = SCREENSHOTS / visual
    if path.exists():
        return path

    # 尝试图表目录
    path = DIAGRAMS_RENDERED / visual
    if path.exists():
        return path

    # 尝试不加后缀搜索
    for ext in ['.png', '.jpg', '.jpeg']:
        path = SCREENSHOTS / (visual + ext)
        if path.exists():
            return path
        path = DIAGRAMS_RENDERED / (visual + ext)
        if path.exists():
            return path

    # 模糊匹配：visual 中包含的关键词
    for directory in [SCREENSHOTS, DIAGRAMS_RENDERED]:
        if directory.exists():
            for f in directory.iterdir():
                if visual.lower() in f.name.lower():
                    return f

    return None


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PPTX 核心工具函数
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def set_gradient_bg(slide, color1: RGBColor, color2: RGBColor, angle=5400000):
    """设置幻灯片渐变背景（默认从上到下）

    正确构建 OOXML 结构：
    <p:cSld>
      <p:bg>
        <p:bgPr>
          <a:gradFill>...</a:gradFill>
          <a:effectLst/>
        </p:bgPr>
      </p:bg>
      <p:spTree>...</p:spTree>
    </p:cSld>
    """
    # 获取 <p:cSld> 元素
    sld = slide._element  # <p:sld>
    cSld = sld.find(qn('p:cSld'))

    # 清除任何已有的 <p:bg> 和误放的 <p:bgPr>
    for old_bg in cSld.findall(qn('p:bg')):
        cSld.remove(old_bg)
    for stray_bgPr in cSld.findall(qn('p:bgPr')):
        cSld.remove(stray_bgPr)

    # 创建 <p:bg> → <p:bgPr> → <a:gradFill> 结构
    bg = etree.Element(qn('p:bg'))
    bgPr = etree.SubElement(bg, qn('p:bgPr'))

    gradFill = etree.SubElement(bgPr, qn('a:gradFill'))
    gradFill.set('rotWithShape', '0')

    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

    gs1 = etree.SubElement(gsLst, qn('a:gs'))
    gs1.set('pos', '0')
    srgb1 = etree.SubElement(gs1, qn('a:srgbClr'))
    srgb1.set('val', str(color1))

    gs2 = etree.SubElement(gsLst, qn('a:gs'))
    gs2.set('pos', '100000')
    srgb2 = etree.SubElement(gs2, qn('a:srgbClr'))
    srgb2.set('val', str(color2))

    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', str(angle))
    lin.set('scaled', '1')

    etree.SubElement(bgPr, qn('a:effectLst'))

    # 插入 <p:bg> 到 <p:spTree> 之前（OOXML 要求的正确位置）
    spTree = cSld.find(qn('p:spTree'))
    if spTree is not None:
        cSld.insert(list(cSld).index(spTree), bg)
    else:
        cSld.insert(0, bg)


def set_solid_bg(slide, color: RGBColor):
    """设置纯色背景（使用与渐变相同的正确 OOXML 结构）"""
    sld = slide._element
    cSld = sld.find(qn('p:cSld'))

    # 清除已有背景
    for old_bg in cSld.findall(qn('p:bg')):
        cSld.remove(old_bg)
    for stray_bgPr in cSld.findall(qn('p:bgPr')):
        cSld.remove(stray_bgPr)

    # 创建 <p:bg> → <p:bgPr> → <a:solidFill>
    bg = etree.Element(qn('p:bg'))
    bgPr = etree.SubElement(bg, qn('p:bgPr'))

    solidFill = etree.SubElement(bgPr, qn('a:solidFill'))
    srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgb.set('val', str(color))

    etree.SubElement(bgPr, qn('a:effectLst'))

    # 插入到 <p:spTree> 之前
    spTree = cSld.find(qn('p:spTree'))
    if spTree is not None:
        cSld.insert(list(cSld).index(spTree), bg)
    else:
        cSld.insert(0, bg)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_BODY, line_spacing=1.2, anchor=None):
    """添加文本框，返回 shape"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.paragraphs[0].alignment = alignment
        # 设置垂直对齐
        txBody = txBox._element.find(qn('p:txBody'))
        if txBody is not None:
            bodyPr = txBody.find(qn('a:bodyPr'))
            if bodyPr is not None:
                anchor_map = {
                    'top': 't', 'middle': 'ctr', 'bottom': 'b'
                }
                bodyPr.set('anchor', anchor_map.get(anchor, 'ctr'))

    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    # East Asian font (WPS needs <a:ea> for Chinese text)
    rPr = run._r.find(qn('a:rPr'))
    if rPr is not None:
        ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', font_name)

    # 行间距
    pPr = p._element.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.SubElement(p._element, qn('a:pPr'))
    lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
    spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
    spcPct.set('val', str(int(line_spacing * 100000)))

    return txBox


def add_multiline_text_box(slide, left, top, width, height, lines,
                           font_size=18, color=TEXT_WHITE, bold=False,
                           alignment=PP_ALIGN.LEFT, font_name=FONT_BODY,
                           line_spacing=1.3, colors=None, sizes=None, bolds=None):
    """添加多行文本框，每行可指定不同颜色/字号/粗细"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = alignment
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(sizes[i] if sizes and i < len(sizes) else font_size)
        run.font.color.rgb = colors[i] if colors and i < len(colors) else color
        run.font.bold = bolds[i] if bolds and i < len(bolds) else bold
        run.font.name = font_name
        # East Asian font
        rPr = run._r.find(qn('a:rPr'))
        if rPr is not None:
            ea = etree.SubElement(rPr, qn('a:ea'))
            ea.set('typeface', font_name)

        # 行间距
        pPr = p._element.find(qn('a:pPr'))
        if pPr is None:
            pPr = etree.SubElement(p._element, qn('a:pPr'))
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))

    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT1, thickness=3):
    """添加装饰性强调线"""
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = color
    line_shape.line.fill.background()
    return line_shape


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None,
                   border_width=0, corner_radius=None):
    """添加矩形（可选圆角、边框）"""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width) if border_width else Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, left, top, diameter, fill_color, border_color=None):
    """添加圆形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_placeholder_box(slide, text):
    """添加占位框（缺少资产时使用）"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3), Inches(2.5), Inches(7), Inches(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x18, 0x1E, 0x30)
    shape.line.color.rgb = ACCENT1
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.color.rgb = TEXT_DIM
    run.font.name = FONT_BODY
    rPr = run._r.find(qn('a:rPr'))
    if rPr is not None:
        ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', FONT_BODY)


def _insert_image_centered(slide, img_path, max_w, max_h, top_offset,
                           use_shadow=False):
    """居中插入图片，可选圆角阴影"""
    try:
        if use_shadow:
            processed = add_rounded_shadow(img_path, radius=16, shadow_offset=6, shadow_blur=12)
            pil_img = Image.open(processed)
            proc_w, proc_h = pil_img.size
            aspect = proc_w / proc_h
            processed.seek(0)
            img_source = processed
        else:
            pil_img = Image.open(img_path)
            proc_w, proc_h = pil_img.size
            aspect = proc_w / proc_h
            img_source = str(img_path)

        if aspect > (max_w / max_h):
            w = max_w
            h = int(w / aspect)
        else:
            h = max_h
            w = int(h * aspect)

        left = (SLIDE_WIDTH - w) // 2
        top = top_offset

        slide.shapes.add_picture(img_source, left, top, w, h)
        return True
    except Exception as e:
        _add_placeholder_box(slide, f"图片加载失败: {e}")
        return False


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 进度条 + 底栏 + Notes
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def add_progress_bar(slide, idx: int, total: int, act: str):
    """顶部3pt薄条进度条"""
    bar_h = PROGRESS_BAR_HEIGHT  # Pt(3)
    bar_y = PROGRESS_BAR_Y       # 0

    # 背景条
    bar_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, bar_y, SLIDE_WIDTH, bar_h
    )
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = PROGRESS_BG
    bar_bg.line.fill.background()

    # 前景条
    progress = (idx + 1) / total if total > 0 else 0
    bar_fg_w = int(SLIDE_WIDTH * progress)
    if bar_fg_w > 0:
        act_color = ACTS.get(act, {}).get('color', ACCENT1)
        bar_fg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, bar_y, bar_fg_w, bar_h
        )
        bar_fg.fill.solid()
        bar_fg.fill.fore_color.rgb = act_color
        bar_fg.line.fill.background()


def add_footer_v3(slide, spec: SlideSpecV3, idx: int, total: int):
    """v3 底栏：幕名 | 页码"""
    footer_y = Inches(7.1)
    footer_h = Inches(0.3)

    # 幕名
    act_info = ACTS.get(spec.act, {})
    act_label = act_info.get('label', spec.act)
    add_text_box(
        slide, Inches(0.5), footer_y, Inches(4), footer_h,
        act_label,
        font_size=FONT_SIZES['footer'], color=TEXT_DIM,
        alignment=PP_ALIGN.LEFT
    )

    # 页码
    add_text_box(
        slide, Inches(5.5), footer_y, Inches(2.5), footer_h,
        f"{idx + 1} / {total}",
        font_size=FONT_SIZES['footer'], color=TEXT_DIM,
        alignment=PP_ALIGN.CENTER
    )

    # 时长
    if spec.duration > 0:
        add_text_box(
            slide, Inches(9.5), footer_y, Inches(3.3), footer_h,
            f"{spec.duration:.1f} min",
            font_size=FONT_SIZES['footer'], color=TEXT_DIM,
            alignment=PP_ALIGN.RIGHT
        )


def add_notes(slide, spec: SlideSpecV3):
    """写入 notes 到备注区"""
    if not spec.notes:
        return
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = spec.notes


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 辅助：获取幕色
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _get_act_color(act: str) -> RGBColor:
    """获取当前幕的强调色"""
    return ACTS.get(act, {}).get('color', ACCENT1)


def _get_act_num(act: str) -> str:
    """获取幕的汉字编号"""
    return ACTS.get(act, {}).get('num', '')


def _dim_color(color: RGBColor, factor=0.2) -> RGBColor:
    """生成暗化版本的颜色（模拟半透明效果）"""
    hex_str = str(color)
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    # 混合背景色 (0x0A, 0x0F, 0x1E)
    bg_r, bg_g, bg_b = 0x0A, 0x0F, 0x1E
    nr = int(bg_r + (r - bg_r) * factor)
    ng = int(bg_g + (g - bg_g) * factor)
    nb = int(bg_b + (b - bg_b) * factor)
    return RGBColor(max(0, min(255, nr)), max(0, min(255, ng)), max(0, min(255, nb)))


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Builder 1: act-title — 幕标题页
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_act_title(prs, spec: SlideSpecV3, idx: int, total: int):
    """
    超深渐变背景
    大号汉字幕号（半透明120pt）
    标题60pt白色居中
    下方强调线
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_SECTION, RGBColor(0x0F, 0x12, 0x2A))

    act_color = _get_act_color(spec.act)
    act_num = _get_act_num(spec.act)

    # 大号汉字幕号（极暗色模拟半透明）
    if act_num:
        dim = _dim_color(act_color, factor=0.15)
        add_text_box(
            slide, Inches(0.5), Inches(0.3), Inches(5), Inches(5),
            act_num,
            font_size=120, color=dim,
            bold=True, font_name=FONT_TITLE
        )

    # 主标题 60pt 居中
    title_text = spec.on_screen or spec.title
    add_text_box(
        slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.5),
        title_text,
        font_size=FONT_SIZES['act-title'], color=TEXT_WHITE, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE,
        line_spacing=1.4
    )

    # 下方强调线
    add_accent_line(
        slide, Inches(5.2), Inches(4.6), Inches(2.9),
        act_color, thickness=4
    )

    return slide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Builder 2: bridge — 过渡页
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_bridge(prs, spec: SlideSpecV3, idx: int, total: int):
    """
    标准深色渐变背景
    居中1句话，28pt
    含"？"= 提问型 → 橙色
    否则 = 回收型 → 青绿色
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    text = spec.on_screen or spec.title
    is_question = '？' in text or '?' in text
    text_color = ACCENT2 if is_question else ACCENT1

    # 居中显示
    add_text_box(
        slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5),
        text,
        font_size=FONT_SIZES['section-title'], color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE,
        line_spacing=1.5
    )

    # 上方细线
    line_color = ACCENT2 if is_question else ACCENT1
    add_accent_line(
        slide, Inches(5.5), Inches(2.2), Inches(2.3),
        line_color, thickness=2
    )

    # 下方细线
    add_accent_line(
        slide, Inches(5.5), Inches(5.2), Inches(2.3),
        line_color, thickness=2
    )

    return slide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Builder 3: roadmap — 导航页
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_roadmap(prs, spec: SlideSpecV3, idx: int, total: int):
    """
    横排显示8个幕的名称
    当前幕高亮（ACCENT1），已过幕dim色，未来幕更dim色
    圆点+连线的时间线样式
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    # 标题
    roadmap_title = spec.on_screen if spec.on_screen and spec.on_screen != "none" else "旅程地图"
    add_text_box(
        slide, Inches(1), Inches(0.8), Inches(11.3), Inches(1.2),
        roadmap_title,
        font_size=FONT_SIZES['section-title'], color=TEXT_WHITE, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE
    )

    current_act = spec.act
    current_idx = ACT_ORDER.index(current_act) if current_act in ACT_ORDER else -1

    num_acts = len(ACT_ORDER)
    total_width = Inches(11.0)
    start_x = Inches(1.16)
    y_center = Inches(3.8)
    dot_diameter = Inches(0.35)
    spacing = total_width // num_acts

    # 连线（横贯所有圆点）
    line_y = y_center + dot_diameter // 2 - Pt(1)
    add_shape_rect(
        slide,
        start_x + dot_diameter // 2,
        line_y,
        spacing * (num_acts - 1),
        Pt(2),
        fill_color=RGBColor(0x2A, 0x35, 0x50)
    )

    for i, act_key in enumerate(ACT_ORDER):
        act_info = ACTS[act_key]
        x = start_x + spacing * i

        if i < current_idx:
            # 已过幕
            dot_color = _dim_color(act_info['color'], factor=0.5)
            label_color = TEXT_DIM
        elif i == current_idx:
            # 当前幕 — 高亮
            dot_color = act_info['color']
            label_color = TEXT_WHITE
        else:
            # 未来幕
            dot_color = TEXT_XDIM
            label_color = TEXT_XDIM

        # 圆点
        add_circle(slide, x, y_center, dot_diameter, dot_color)

        # 幕名标签
        label = act_info['label']
        label_w = Inches(1.3)
        label_x = x - (label_w - dot_diameter) // 2
        add_text_box(
            slide, label_x, y_center + Inches(0.55), label_w, Inches(0.5),
            label,
            font_size=13 if i == current_idx else 11,
            color=label_color, bold=(i == current_idx),
            alignment=PP_ALIGN.CENTER, font_name=FONT_BODY
        )

        # 当前幕加外圈
        if i == current_idx:
            ring = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x - Inches(0.06),
                y_center - Inches(0.06),
                dot_diameter + Inches(0.12),
                dot_diameter + Inches(0.12)
            )
            ring.fill.background()
            ring.line.color.rgb = act_info['color']
            ring.line.width = Pt(2)

    return slide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Builder 4: key-phrase — 核心短语页
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_key_phrase(prs, spec: SlideSpecV3, idx: int, total: int):
    """
    标准深色渐变背景
    on-screen文字44pt白色居中
    上下有装饰性细线（ACCENT1/ACCENT2交替）
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    text = spec.on_screen or spec.title
    act_color = _get_act_color(spec.act)

    # 交替装饰色：偶数页用 ACCENT1，奇数页用 ACCENT2
    deco_color = ACCENT1 if idx % 2 == 0 else ACCENT2

    # 主文字 44pt 居中（先于装饰线添加，确保 WPS 渲染）
    add_text_box(
        slide, Inches(1.2), Inches(2.3), Inches(10.9), Inches(3.0),
        text,
        font_size=FONT_SIZES['key-phrase'], color=TEXT_WHITE, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE,
        line_spacing=1.4
    )

    # 上方装饰线
    add_accent_line(
        slide, Inches(4.5), Inches(2.0), Inches(4.3),
        deco_color, thickness=2
    )

    # 下方装饰线
    add_accent_line(
        slide, Inches(4.5), Inches(5.4), Inches(4.3),
        deco_color, thickness=2
    )

    return slide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Builder 5: full-visual — 全图页
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_full_visual(prs, spec: SlideSpecV3, idx: int, total: int):
    """
    有visual：图片居中大面积 + 顶部标题28pt
    无visual：大号标题居中 + caption
    on-screen 通常是 "标题\\ncaption"
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    img_path = resolve_visual(spec)
    on_lines = spec.on_screen.split('\n') if spec.on_screen else [spec.title]
    title_text = on_lines[0].strip() if on_lines else spec.title
    caption_text = on_lines[1].strip() if len(on_lines) > 1 else ""

    if img_path and img_path.exists():
        # 顶部标题
        add_text_box(
            slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.8),
            title_text,
            font_size=FONT_SIZES['section-title'], color=ACCENT1, bold=True,
            font_name=FONT_TITLE
        )
        add_accent_line(
            slide, Inches(0.6), Inches(1.05), Inches(3),
            _get_act_color(spec.act), thickness=3
        )

        # 图片居中
        _insert_image_centered(
            slide, img_path,
            max_w=Inches(11.5), max_h=Inches(5.3),
            top_offset=Inches(1.3),
            use_shadow=False
        )

        # caption
        if caption_text:
            add_text_box(
                slide, Inches(1), Inches(6.6), Inches(11.3), Inches(0.4),
                caption_text,
                font_size=FONT_SIZES['caption'], color=TEXT_DIM,
                alignment=PP_ALIGN.CENTER
            )
    else:
        # 无图 — 大号标题居中
        add_text_box(
            slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.5),
            title_text,
            font_size=FONT_SIZES['key-phrase'], color=TEXT_WHITE, bold=True,
            alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE,
            line_spacing=1.4
        )
        if caption_text:
            add_text_box(
                slide, Inches(2), Inches(4.8), Inches(9.3), Inches(1),
                caption_text,
                font_size=FONT_SIZES['caption'], color=TEXT_DIM,
                alignment=PP_ALIGN.CENTER
            )

    return slide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Builder 6: evidence — 证据/截图页
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_evidence(prs, spec: SlideSpecV3, idx: int, total: int):
    """
    截图占60%页面（居中偏右），带圆角阴影
    左上标题28pt
    on-screen中除标题外的行作为标注callout
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    img_path = resolve_visual(spec)
    on_lines = spec.on_screen.split('\n') if spec.on_screen else [spec.title]
    title_text = on_lines[0].strip() if on_lines else spec.title
    callout_lines = [l.strip() for l in on_lines[1:] if l.strip()]

    # 标题 — 左上
    add_text_box(
        slide, Inches(0.6), Inches(0.3), Inches(5.8), Inches(0.9),
        title_text,
        font_size=FONT_SIZES['section-title'], color=ACCENT1, bold=True,
        font_name=FONT_TITLE
    )
    add_accent_line(
        slide, Inches(0.6), Inches(1.2), Inches(2.5),
        _get_act_color(spec.act), thickness=3
    )

    # Callout 标注 — 左侧
    if callout_lines:
        for ci, cline in enumerate(callout_lines):
            add_text_box(
                slide, Inches(0.6), Inches(1.7 + ci * 0.65),
                Inches(5.2), Inches(0.55),
                cline,
                font_size=15, color=TEXT_BODY, bold=False,
                font_name=FONT_BODY, line_spacing=1.3
            )

    # 截图 — 右侧，带圆角阴影
    if img_path and img_path.exists():
        try:
            processed = add_rounded_shadow(img_path, radius=16, shadow_offset=6, shadow_blur=12)
            pil_img = Image.open(processed)
            proc_w, proc_h = pil_img.size
            aspect = proc_w / proc_h
            processed.seek(0)

            max_w = Inches(6.5)
            max_h = Inches(5.5)

            if aspect > (max_w / max_h):
                w = max_w
                h = int(w / aspect)
            else:
                h = max_h
                w = int(h * aspect)

            left = SLIDE_WIDTH - w - Inches(0.4)
            top = Inches(1.0)
            slide.shapes.add_picture(processed, left, top, w, h)
        except Exception as e:
            add_text_box(
                slide, Inches(7), Inches(2.5), Inches(5.5), Inches(2),
                f"[图片: {spec.visual}]",
                font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER
            )
    else:
        if spec.visual and spec.visual != "none":
            add_text_box(
                slide, Inches(7), Inches(3), Inches(5), Inches(1),
                f"[缺失: {spec.visual}]",
                font_size=12, color=RGBColor(0xFF, 0x44, 0x44),
                alignment=PP_ALIGN.CENTER
            )

    return slide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Builder 7: split-compare — 左右对比页
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_split_compare(prs, spec: SlideSpecV3, idx: int, total: int):
    """
    页面左右对分
    中间 VS 分隔符
    on-screen 用 "|" 或 "vs" 分割左右
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    text = spec.on_screen or spec.title

    # 解析左右
    if '|' in text:
        parts = text.split('|', 1)
    elif ' vs ' in text.lower():
        parts = re.split(r'\s+vs\s+', text, maxsplit=1, flags=re.IGNORECASE)
    elif '\n' in text:
        lines = text.split('\n')
        mid = len(lines) // 2
        parts = ['\n'.join(lines[:mid]), '\n'.join(lines[mid:])]
    else:
        parts = [text, ""]

    left_text = parts[0].strip()
    right_text = parts[1].strip() if len(parts) > 1 else ""

    half_w = SLIDE_WIDTH // 2

    # 左半暗底
    add_shape_rect(
        slide, 0, Inches(0.6), half_w - Inches(0.3), Inches(6.3),
        fill_color=RGBColor(0x0D, 0x12, 0x22),
        corner_radius=True
    )

    # 右半暗底
    add_shape_rect(
        slide, half_w + Inches(0.3), Inches(0.6), half_w - Inches(0.3), Inches(6.3),
        fill_color=RGBColor(0x12, 0x0D, 0x22),
        corner_radius=True
    )

    # 左侧文字
    add_text_box(
        slide, Inches(0.8), Inches(1.5), Inches(5.3), Inches(4.5),
        left_text,
        font_size=22, color=ACCENT1, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE,
        line_spacing=1.5
    )

    # 右侧文字
    add_text_box(
        slide, Inches(7.2), Inches(1.5), Inches(5.3), Inches(4.5),
        right_text,
        font_size=22, color=ACCENT2, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE,
        line_spacing=1.5
    )

    # 中间 VS 分隔符
    center_x = half_w - Inches(0.4)
    # VS 圆
    vs_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        center_x, Inches(3.1), Inches(0.8), Inches(0.8)
    )
    vs_circle.fill.solid()
    vs_circle.fill.fore_color.rgb = RGBColor(0x1A, 0x1E, 0x30)
    vs_circle.line.color.rgb = BORDER_COLOR
    vs_circle.line.width = Pt(2)

    # VS 文字
    add_text_box(
        slide, center_x - Inches(0.1), Inches(3.15), Inches(1.0), Inches(0.7),
        "VS",
        font_size=16, color=TEXT_WHITE, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE
    )

    # 竖线装饰（VS 上下）
    vline_x = half_w - Pt(1)
    add_shape_rect(slide, vline_x, Inches(0.8), Pt(2), Inches(2.2),
                   fill_color=BORDER_COLOR)
    add_shape_rect(slide, vline_x, Inches(4.2), Pt(2), Inches(2.5),
                   fill_color=BORDER_COLOR)

    # 如果有 visual，尝试作为背景的某一侧
    img_path = resolve_visual(spec)
    if img_path and img_path.exists():
        try:
            processed = add_rounded_shadow(img_path, radius=12, shadow_offset=4, shadow_blur=8)
            pil_img = Image.open(processed)
            aspect = pil_img.width / pil_img.height
            processed.seek(0)
            max_w = Inches(5.5)
            max_h = Inches(3.0)
            if aspect > (max_w / max_h):
                w = max_w
                h = int(w / aspect)
            else:
                h = max_h
                w = int(h * aspect)
            # 放在右下角
            slide.shapes.add_picture(
                processed,
                SLIDE_WIDTH - w - Inches(0.5),
                Inches(4.0),
                w, h
            )
        except Exception:
            pass

    return slide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Builder 8: data-callout — 大数字页
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_data_callout(prs, spec: SlideSpecV3, idx: int, total: int):
    """
    大数字72pt+，居中偏上
    说明文字18pt居中偏下
    on-screen 格式："数字\\n说明"
    数字用ACCENT1/ACCENT2色
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    text = spec.on_screen or spec.title
    lines = text.split('\n')
    number_text = lines[0].strip() if lines else text
    desc_text = '\n'.join(l.strip() for l in lines[1:]).strip() if len(lines) > 1 else ""

    # 交替色
    num_color = ACCENT1 if idx % 2 == 0 else ACCENT2
    act_color = _get_act_color(spec.act)

    # 大数字 72pt
    add_text_box(
        slide, Inches(1), Inches(1.5), Inches(11.3), Inches(3.0),
        number_text,
        font_size=FONT_SIZES['data-number'], color=num_color, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE,
        line_spacing=1.1
    )

    # 数字下方强调线
    add_accent_line(
        slide, Inches(5), Inches(4.5), Inches(3.3),
        act_color, thickness=3
    )

    # 说明文字
    if desc_text:
        add_text_box(
            slide, Inches(2), Inches(4.9), Inches(9.3), Inches(1.8),
            desc_text,
            font_size=18, color=TEXT_BODY, bold=False,
            alignment=PP_ALIGN.CENTER, font_name=FONT_BODY,
            line_spacing=1.4
        )

    return slide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Builder 9: quote-impact — 金句页
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_quote_impact(prs, spec: SlideSpecV3, idx: int, total: int):
    """
    深色渐变（偏紫调）
    装饰大引号（左上+右下）
    金句36pt白色居中
    底部细线装饰
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, QUOTE_BG1, QUOTE_BG2, angle=2700000)

    quote_text = spec.on_screen or spec.title

    # 根据文本长度调整字号
    q_len = len(quote_text)
    if q_len > 60:
        q_size = 26
    elif q_len > 40:
        q_size = 30
    elif q_len > 25:
        q_size = 34
    else:
        q_size = FONT_SIZES['quote']

    # 装饰大引号 — 左上
    add_text_box(
        slide, Inches(1.0), Inches(0.8), Inches(2), Inches(2),
        "\u201C",
        font_size=120, color=RGBColor(0x00, 0xD4, 0xAA),
        bold=True, font_name="Georgia"
    )

    # 装饰大引号 — 右下
    add_text_box(
        slide, Inches(10.3), Inches(4.5), Inches(2), Inches(2),
        "\u201D",
        font_size=120, color=RGBColor(0x00, 0x80, 0x66),
        bold=True, font_name="Georgia"
    )

    # 金句主体
    add_text_box(
        slide, Inches(2), Inches(2.0), Inches(9.3), Inches(3.5),
        quote_text,
        font_size=q_size, color=TEXT_WHITE, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE,
        line_spacing=1.6
    )

    # 底部细线装饰
    add_accent_line(
        slide, Inches(4), Inches(5.8), Inches(5.3),
        ACCENT1, thickness=2
    )

    return slide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Builder 10: sequence — 序列/流程页
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_sequence(prs, spec: SlideSpecV3, idx: int, total: int):
    """
    横排3-4格，每格等宽
    从 on-screen 解析标签（按 \\n 或 → 分割）
    格之间有箭头连接
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    text = spec.on_screen or spec.title

    # 解析标签
    if '\u2192' in text:  # →
        items = [s.strip() for s in text.split('\u2192') if s.strip()]
    elif '->' in text:
        items = [s.strip() for s in text.split('->') if s.strip()]
    elif '\n' in text:
        items = [s.strip() for s in text.split('\n') if s.strip()]
    elif '|' in text:
        items = [s.strip() for s in text.split('|') if s.strip()]
    else:
        items = [text]

    # 如果第一个是标题行（没有箭头但存在后续项）
    title_text = spec.title
    if len(items) >= 2:
        # 检查第一项是否像标题（比其他项长很多或含冒号）
        first = items[0]
        if ':' in first or '：' in first:
            title_text = first
            items = items[1:]

    # 标题
    add_text_box(
        slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
        title_text,
        font_size=FONT_SIZES['section-title'], color=ACCENT1, bold=True,
        font_name=FONT_TITLE
    )
    add_accent_line(
        slide, Inches(0.6), Inches(1.1), Inches(3),
        _get_act_color(spec.act), thickness=3
    )

    # 布局序列格
    num_items = max(len(items), 1)
    num_items = min(num_items, 6)  # 最多6格
    items = items[:num_items]

    margin_x = Inches(0.8)
    total_area_w = SLIDE_WIDTH - 2 * margin_x
    gap = Inches(0.6)
    # 箭头宽度
    arrow_w = Inches(0.4)
    # 可用于格子的总宽度 = total_area_w - (num-1) * (gap + arrow_w)
    usable_w = total_area_w - (num_items - 1) * (gap + arrow_w)
    cell_w = usable_w // num_items if num_items > 0 else total_area_w
    cell_h = Inches(3.2)
    y_top = Inches(2.0)
    act_color = _get_act_color(spec.act)

    for i, item in enumerate(items):
        x = margin_x + i * (cell_w + gap + arrow_w)

        # 格子背景
        add_shape_rect(
            slide, x, y_top, cell_w, cell_h,
            fill_color=RGBColor(0x10, 0x15, 0x28),
            border_color=_dim_color(act_color, factor=0.4),
            border_width=1,
            corner_radius=True
        )

        # 序号圆点
        num_d = Inches(0.35)
        add_circle(
            slide, x + (cell_w - num_d) // 2, y_top + Inches(0.2),
            num_d, act_color
        )
        add_text_box(
            slide, x + (cell_w - num_d) // 2, y_top + Inches(0.2),
            num_d, num_d,
            str(i + 1),
            font_size=14, color=RGBColor(0x0A, 0x0F, 0x1E), bold=True,
            alignment=PP_ALIGN.CENTER, font_name=FONT_BODY
        )

        # 标签文字
        add_text_box(
            slide, x + Inches(0.15), y_top + Inches(0.75),
            cell_w - Inches(0.3), cell_h - Inches(1.0),
            item,
            font_size=15, color=TEXT_BODY, bold=False,
            alignment=PP_ALIGN.CENTER, font_name=FONT_BODY,
            line_spacing=1.3
        )

        # 箭头（除最后一格外）
        if i < num_items - 1:
            arrow_x = x + cell_w + gap // 2
            arrow_y = y_top + cell_h // 2 - Inches(0.15)
            add_text_box(
                slide, arrow_x, arrow_y, arrow_w, Inches(0.3),
                "\u25B6",  # ▶
                font_size=18, color=act_color, bold=False,
                alignment=PP_ALIGN.CENTER
            )

    # 如果有 visual，放在底部小图
    img_path = resolve_visual(spec)
    if img_path and img_path.exists():
        try:
            _insert_image_centered(
                slide, img_path,
                max_w=Inches(4), max_h=Inches(1.5),
                top_offset=Inches(5.5),
                use_shadow=False
            )
        except Exception:
            pass

    return slide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Builder 11: checklist — 检查单页
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_checklist(prs, spec: SlideSpecV3, idx: int, total: int):
    """
    飞行员检查单样式
    每行一个 ✓ + 条目
    白色文字 + 绿色对勾
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    text = spec.on_screen or spec.title
    lines = [l.strip() for l in text.split('\n') if l.strip()]

    # 第一行可能是标题
    title_text = spec.title
    check_items = lines
    if lines and (':' in lines[0] or '：' in lines[0]):
        title_text = lines[0]
        check_items = lines[1:]
    elif not check_items:
        check_items = [text]

    # 标题
    add_text_box(
        slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.9),
        title_text,
        font_size=FONT_SIZES['section-title'], color=ACCENT1, bold=True,
        font_name=FONT_TITLE
    )
    add_accent_line(
        slide, Inches(0.8), Inches(1.2), Inches(3),
        _get_act_color(spec.act), thickness=3
    )

    # 检查单容器背景
    list_x = Inches(1.5)
    list_y = Inches(1.6)
    list_w = Inches(10.3)
    item_h = Inches(0.65)

    check_color = RGBColor(0x00, 0xD4, 0xAA)  # 绿色对勾

    for i, item in enumerate(check_items):
        y = list_y + i * item_h

        # 行背景（交替微妙色差）
        row_bg_color = RGBColor(0x0D, 0x12, 0x22) if i % 2 == 0 else RGBColor(0x0F, 0x14, 0x26)
        add_shape_rect(
            slide, list_x, y, list_w, item_h - Inches(0.05),
            fill_color=row_bg_color,
            corner_radius=True
        )

        # 对勾
        add_text_box(
            slide, list_x + Inches(0.2), y + Inches(0.08),
            Inches(0.5), item_h - Inches(0.1),
            "\u2713",  # ✓
            font_size=20, color=check_color, bold=True,
            alignment=PP_ALIGN.CENTER, font_name=FONT_BODY
        )

        # 条目文字
        # 清理已有的前缀符号
        clean_item = item.lstrip('\u2713\u2714\u2715\u2716 -\u2022\u25b8\u25c6\u25cf')
        clean_item = clean_item.strip()
        add_text_box(
            slide, list_x + Inches(0.7), y + Inches(0.08),
            list_w - Inches(1.0), item_h - Inches(0.1),
            clean_item,
            font_size=17, color=TEXT_WHITE, bold=False,
            font_name=FONT_BODY
        )

        # 右侧虚线（模拟航空检查单风格）
        dots_text = "\u00B7" * 20  # middot
        add_text_box(
            slide, list_x + Inches(0.7), y + Inches(0.35),
            list_w - Inches(1.0), Inches(0.2),
            dots_text,
            font_size=8, color=RGBColor(0x25, 0x2A, 0x3A),
            font_name=FONT_BODY
        )

    return slide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Builder 路由表
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
BUILDERS = {
    'act-title':     build_act_title,
    'bridge':        build_bridge,
    'roadmap':       build_roadmap,
    'key-phrase':    build_key_phrase,
    'full-visual':   build_full_visual,
    'evidence':      build_evidence,
    'split-compare': build_split_compare,
    'data-callout':  build_data_callout,
    'quote-impact':  build_quote_impact,
    'sequence':      build_sequence,
    'checklist':     build_checklist,
}


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 主函数
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def main():
    print("=" * 65)
    print("  PPTX Build v3 — AI 科普演讲（11类页面 + 导航系统）")
    print("=" * 65)

    # 1. 解析 spec 文件
    specs = []
    for f in SPEC_FILES:
        if f.exists():
            print(f"\n  解析: {f.name}")
            parsed = parse_v3_spec(f)
            print(f"    -> {len(parsed)} 页")
            specs.extend(parsed)
        else:
            print(f"\n  [跳过] 文件不存在: {f.name}")

    if not specs:
        print("\n  错误：没有找到任何 spec 文件！")
        return

    specs.sort(key=lambda s: s.page_num)
    total = len(specs)

    # 2. 类型统计
    type_counts = {}
    for s in specs:
        t = s.slide_type or "unknown"
        type_counts[t] = type_counts.get(t, 0) + 1
    print(f"\n  总计: {total} 页")
    print("  页面类型分布:")
    for t, c in sorted(type_counts.items()):
        print(f"    {t:16s}: {c}")

    # 3. 检查 visual 资产
    missing_visuals = []
    for s in specs:
        if s.visual and s.visual != "none":
            path = resolve_visual(s)
            if path is None:
                missing_visuals.append((s.page_num, s.visual))
    if missing_visuals:
        print(f"\n  警告: {len(missing_visuals)} 个 visual 缺失:")
        for pn, vn in missing_visuals[:10]:
            print(f"    P{pn}: {vn}")
        if len(missing_visuals) > 10:
            print(f"    ... 还有 {len(missing_visuals) - 10} 个")

    # 4. 构建 PPTX
    print(f"\n  开始构建 PPTX ({total} 页)...")
    print("-" * 65)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    stats = {}
    for idx_i, spec in enumerate(specs):
        builder = BUILDERS.get(spec.slide_type, build_key_phrase)
        stype = spec.slide_type if spec.slide_type in BUILDERS else f"key-phrase(fallback:{spec.slide_type})"

        try:
            slide = builder(prs, spec, idx_i, total)
            add_progress_bar(slide, idx_i, total, spec.act)
            add_footer_v3(slide, spec, idx_i, total)
            add_notes(slide, spec)
        except Exception as e:
            # 出错时使用 key-phrase 兜底
            print(f"  [ERROR] P{spec.page_num} ({stype}): {e}")
            slide = build_key_phrase(prs, spec, idx_i, total)
            add_progress_bar(slide, idx_i, total, spec.act)
            add_footer_v3(slide, spec, idx_i, total)
            add_notes(slide, spec)
            stype = f"key-phrase(error:{spec.slide_type})"

        stats[stype] = stats.get(stype, 0) + 1
        title_preview = (spec.on_screen or spec.title)[:35].replace('\n', ' ')
        print(f"  [{idx_i+1:3d}/{total}] P{spec.page_num:3d} [{stype:16s}] {title_preview}")

    print("-" * 65)

    # 5. 保存
    prs.save(str(OUTPUT_PPTX))
    size_mb = OUTPUT_PPTX.stat().st_size / (1024 * 1024)

    print(f"\n  已保存: {OUTPUT_PPTX}")
    print(f"  文件大小: {size_mb:.1f} MB")
    print(f"  总页数: {total}")
    print(f"\n  构建统计:")
    for k, v in sorted(stats.items()):
        print(f"    {k:24s}: {v}")

    # 6. 质量检查
    print(f"\n  质量检查:")
    print(f"    visual 缺失: {len(missing_visuals)}")
    print(f"    文件 < 30MB: {'OK' if size_mb < 30 else 'WARN'} ({size_mb:.1f} MB)")
    print(f"    11 builder 覆盖: {len(BUILDERS)} 种类型")
    print("=" * 65)


if __name__ == "__main__":
    main()
