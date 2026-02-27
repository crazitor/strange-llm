#!/usr/bin/env python3
"""
build_pptx_v2.py — 高质量 PPTX 构建脚本
从 slides_spec_v2 的 markdown 规格文件构建专业级演示文稿

特性：
- 渐变背景（深海蓝→暗紫）
- 截图圆角阴影处理
- 金句卡渐变 + 装饰引号
- 幕标题装饰 + 章节编号
- 底栏页码 + 章节进度条
- 图表占位替换为实际 PNG
"""

import re
import os
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional
from copy import deepcopy
from io import BytesIO
import json

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree
from PIL import Image, ImageDraw, ImageFilter

# ─── 路径配置 ───────────────────────────────────
PROJECT = Path("/Users/mik/strange LLM")
OUTPUT = PROJECT / "output"
SCREENSHOTS = OUTPUT / "screenshots"
DIAGRAMS_RENDERED = PROJECT / "assets" / "diagrams" / "rendered"
DIAGRAM_INDEX = OUTPUT / "diagram_index.md"
SPEC_FILES = [
    OUTPUT / "slides_spec_v4.md",
]
OUTPUT_PPTX = OUTPUT / "presentation_v4.pptx"

# ─── 设计参数 ───────────────────────────────────
SLIDE_WIDTH = Inches(13.333)   # 16:9 宽屏
SLIDE_HEIGHT = Inches(7.5)

# 配色方案
BG_PRIMARY    = RGBColor(0x0F, 0x17, 0x29)  # 深海蓝
BG_SECONDARY  = RGBColor(0x1A, 0x1A, 0x2E)  # 渐变终点
BG_SECTION    = RGBColor(0x08, 0x0E, 0x1F)  # 幕标题背景（更深）
ACCENT1       = RGBColor(0x00, 0xD4, 0xAA)  # 明亮青绿
ACCENT2       = RGBColor(0xFF, 0x6B, 0x35)  # 活力橙
TEXT_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_BODY      = RGBColor(0xF0, 0xF0, 0xF0)
TEXT_DIM       = RGBColor(0x9C, 0xA3, 0xAF)
QUOTE_BG1      = RGBColor(0x0A, 0x0E, 0x1A)  # 金句渐变起点
QUOTE_BG2      = RGBColor(0x1A, 0x10, 0x30)  # 金句渐变终点（紫调）
PROGRESS_BG    = RGBColor(0x1E, 0x25, 0x3A)  # 进度条背景
BORDER_COLOR   = RGBColor(0x2A, 0x35, 0x50)  # 微妙边框

FONT_TITLE = "PingFang SC"
FONT_BODY_NAME = "PingFang SC"
FONT_FALLBACK = "Microsoft YaHei"

# 章节信息
SECTIONS = {
    "序幕": {"num": "序", "color": ACCENT1, "pages": (1, 3)},
    "第一幕": {"num": "壹", "color": ACCENT1, "pages": (4, 13)},
    "第二幕": {"num": "贰", "color": ACCENT2, "pages": (14, 19)},
    "第三幕": {"num": "叁", "color": ACCENT1, "pages": (20, 29)},
    "第四幕": {"num": "肆", "color": ACCENT2, "pages": (30, 37)},
    "第五幕": {"num": "伍", "color": ACCENT1, "pages": (38, 46)},
    "第六幕": {"num": "陆", "color": ACCENT2, "pages": (47, 50)},
    "第七幕": {"num": "柒", "color": ACCENT1, "pages": (51, 55)},
}

# 图表页码→文件映射
DIAGRAM_MAP = {
    "2":   "patch_tower_layers.png",
    "13":  "patch_tower_full.png",
    "18":  "patch_tower_stoploss.png",
    "22":  "bamboo_growth_curve.png",
    "26":  "capability_inversion.png",
    "26b": "ood_boundary.png",
    "28":  "patch_tower_ceiling.png",
    "33":  "solution_space.png",
    "40":  "climbing_vs_helicopter.png",
    "45":  "ai_vs_human_cost.png",
    "51":  "three_paths.png",
}

# 自定义要点符号
BULLET_SYMBOLS = ["▸", "◆", "●", "▹"]


# ─── 数据结构 ───────────────────────────────────
@dataclass
class SlideSpec:
    page_num: str = "0"
    title: str = ""
    section: str = ""
    duration: str = ""
    cumulative: str = ""
    ppt_title: str = ""
    ppt_subtitle_points: list = field(default_factory=list)
    image_type: str = ""
    image_value: str = ""
    layout: str = ""
    table_data: list = field(default_factory=list)
    speech: str = ""
    speaker_note: str = ""


# ─── Markdown 解析（与原脚本相同逻辑） ──────────
def parse_spec_file(filepath: Path) -> list[SlideSpec]:
    text = filepath.read_text(encoding="utf-8")
    slides = []
    page_pattern = re.compile(r'^### 页 (\d+[A-Za-z]?): (.+)$', re.MULTILINE)
    matches = list(page_pattern.finditer(text))

    for i, match in enumerate(matches):
        start = match.start()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        block = text[start:end]

        spec = SlideSpec()
        spec.page_num = match.group(1)
        spec.title = match.group(2).strip()

        m = re.search(r'\*\*章节\*\*:\s*(.+)', block)
        if m:
            spec.section = m.group(1).strip()

        m = re.search(r'\*\*时间\*\*:\s*(.+)', block)
        if m:
            spec.duration = m.group(1).strip()

        m = re.search(r'\*\*累计时间\*\*:\s*(.+)', block)
        if m:
            spec.cumulative = m.group(1).strip()

        ppt_section = re.search(
            r'\*\*PPT显示内容\*\*:\s*\n(.*?)(?=\n\*\*演讲词\*\*)',
            block, re.DOTALL
        )
        if ppt_section:
            ppt_block = ppt_section.group(1)

            m = re.search(r'- 标题:\s*(.+)', ppt_block)
            if m:
                spec.ppt_title = m.group(1).strip()

            points = []
            in_points = False
            for line in ppt_block.split('\n'):
                if '副标题/要点' in line or '副标题' in line:
                    in_points = True
                    after = line.split(':', 1)
                    if len(after) > 1 and after[1].strip():
                        points.append(after[1].strip())
                    continue
                if in_points and '配图' in line:
                    in_points = False
                elif in_points and line.strip().startswith('-'):
                    point = line.strip().lstrip('-').strip()
                    points.append(point)
                elif in_points and line.strip().startswith('·'):
                    point = line.strip().lstrip('·').strip()
                    points.append(point)
                elif in_points and not line.strip():
                    continue
            spec.ppt_subtitle_points = points

            # 解析布局类型
            layout_m = re.search(r'- 布局:\s*(.+)', ppt_block)
            if layout_m:
                spec.layout = layout_m.group(1).strip()

            # 解析表格数据
            table_lines = []
            in_table = False
            for line in ppt_block.split('\n'):
                if line.strip().startswith('|') and '---' not in line:
                    table_lines.append([c.strip() for c in line.strip().strip('|').split('|')])
                    in_table = True
                elif in_table and not line.strip().startswith('|'):
                    in_table = False
            if table_lines:
                spec.table_data = table_lines

            m = re.search(r'- 配图:\s*(.+)', ppt_block)
            if not m:
                m = re.search(r'配图:\s*(.+)', ppt_block)
            if m:
                img_val = m.group(1).strip()
                if 'diagram-placeholder' in img_val.lower():
                    spec.image_type = "diagram-placeholder"
                    desc = re.search(r'[—–-]\s*(.+?)[\]】]', img_val)
                    spec.image_value = desc.group(1) if desc else img_val
                elif 'quote-card' in img_val.lower() or 'quote_card' in img_val.lower():
                    spec.image_type = "quote-card"
                    desc = re.search(r'[—–-]\s*(.+?)[\]】]', img_val)
                    spec.image_value = desc.group(1) if desc else img_val
                elif 'text-only' in img_val.lower() or 'text_only' in img_val.lower():
                    spec.image_type = "text-only"
                    desc = re.search(r'[—–-]\s*(.+?)[\]】]', img_val)
                    spec.image_value = desc.group(1) if desc else img_val
                elif (img_val.startswith('video') or img_val.startswith('v5_slide')) and (img_val.endswith('.jpg') or img_val.endswith('.png')):
                    spec.image_type = "screenshot"
                    spec.image_value = img_val
                elif '.jpg' in img_val or '.png' in img_val:
                    spec.image_type = "screenshot"
                    fname = re.search(r'((?:video\d+_t\d+_.+?|v5_slide\d+_\d+)\.(jpg|png))', img_val)
                    spec.image_value = fname.group(1) if fname else img_val
                else:
                    spec.image_type = "text-only"
                    spec.image_value = img_val

        speech_match = re.search(
            r'\*\*演讲词\*\*:\s*\n(.*?)(?=\n\*\*讲者提醒\*\*)',
            block, re.DOTALL
        )
        if speech_match:
            speech_lines = speech_match.group(1).strip()
            speech_lines = re.sub(r'^>\s?', '', speech_lines, flags=re.MULTILINE)
            spec.speech = speech_lines.strip()

        m = re.search(r'\*\*讲者提醒\*\*:\s*(.+?)(?:\n---|\n\n|\Z)', block, re.DOTALL)
        if m:
            spec.speaker_note = m.group(1).strip()

        slides.append(spec)

    return slides


# ─── 图片处理工具 ─────────────────────────────────
def add_rounded_shadow(img_path: Path, radius=20, shadow_offset=8, shadow_blur=15) -> BytesIO:
    """给图片添加圆角和阴影效果，返回处理后的 BytesIO"""
    img = Image.open(img_path).convert("RGBA")

    # 缩放到合理尺寸（最大宽 800px 节省 PPTX 体积）
    max_w = 900
    if img.width > max_w:
        ratio = max_w / img.width
        img = img.resize((max_w, int(img.height * ratio)), Image.LANCZOS)

    # 创建圆角蒙版
    mask = Image.new("L", img.size, 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle([(0, 0), img.size], radius=radius, fill=255)

    # 应用圆角
    rounded = Image.new("RGBA", img.size, (0, 0, 0, 0))
    rounded.paste(img, mask=mask)

    # 创建阴影
    canvas_w = img.width + shadow_offset * 2 + shadow_blur * 2
    canvas_h = img.height + shadow_offset * 2 + shadow_blur * 2
    canvas = Image.new("RGBA", (canvas_w, canvas_h), (0, 0, 0, 0))

    # 阴影层
    shadow = Image.new("RGBA", img.size, (0, 0, 0, 80))
    shadow_masked = Image.new("RGBA", img.size, (0, 0, 0, 0))
    shadow_masked.paste(shadow, mask=mask)

    shadow_canvas = Image.new("RGBA", (canvas_w, canvas_h), (0, 0, 0, 0))
    shadow_canvas.paste(shadow_masked,
                        (shadow_blur + shadow_offset, shadow_blur + shadow_offset))
    shadow_canvas = shadow_canvas.filter(ImageFilter.GaussianBlur(shadow_blur))

    canvas = Image.alpha_composite(canvas, shadow_canvas)
    canvas.paste(rounded, (shadow_blur, shadow_blur), rounded)

    # 添加细边框
    border_draw = ImageDraw.Draw(canvas)
    border_draw.rounded_rectangle(
        [(shadow_blur, shadow_blur),
         (shadow_blur + img.width - 1, shadow_blur + img.height - 1)],
        radius=radius,
        outline=(74, 201, 176, 60),  # ACCENT1 半透明
        width=2
    )

    buf = BytesIO()
    canvas.save(buf, format="PNG")
    buf.seek(0)
    return buf


# ─── PPTX 核心构建函数 ───────────────────────────
def set_gradient_bg(slide, color1: RGBColor, color2: RGBColor, angle=5400000):
    """设置幻灯片渐变背景（从上到下）"""
    cSld = slide._element.find(qn('p:cSld'))
    # 移除已有 p:bg
    old_bg = cSld.find(qn('p:bg'))
    if old_bg is not None:
        cSld.remove(old_bg)

    # 创建 p:bg > p:bgPr > a:gradFill 结构，插入到 cSld 的第一个子元素位置
    bg = etree.SubElement(cSld, qn('p:bg'))
    cSld.remove(bg)
    cSld.insert(0, bg)

    bgPr = etree.SubElement(bg, qn('p:bgPr'))

    gradFill = etree.SubElement(bgPr, qn('a:gradFill'))
    gradFill.set('rotWithShape', '0')

    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

    # Stop 1: color1 at 0%
    gs1 = etree.SubElement(gsLst, qn('a:gs'))
    gs1.set('pos', '0')
    srgb1 = etree.SubElement(gs1, qn('a:srgbClr'))
    srgb1.set('val', str(color1))

    # Stop 2: color2 at 100%
    gs2 = etree.SubElement(gsLst, qn('a:gs'))
    gs2.set('pos', '100000')
    srgb2 = etree.SubElement(gs2, qn('a:srgbClr'))
    srgb2.set('val', str(color2))

    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', str(angle))   # 5400000 = 上到下
    lin.set('scaled', '1')

    # 效果
    etree.SubElement(bgPr, qn('a:effectLst'))


def set_solid_bg(slide, color: RGBColor):
    """设置纯色背景"""
    cSld = slide._element.find(qn('p:cSld'))
    old_bg = cSld.find(qn('p:bg'))
    if old_bg is not None:
        cSld.remove(old_bg)

    bg = etree.SubElement(cSld, qn('p:bg'))
    cSld.remove(bg)
    cSld.insert(0, bg)

    bgPr = etree.SubElement(bg, qn('p:bgPr'))
    solidFill = etree.SubElement(bgPr, qn('a:solidFill'))
    srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgb.set('val', str(color))
    etree.SubElement(bgPr, qn('a:effectLst'))


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_BODY_NAME, line_spacing=1.2):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    # 设置行间距
    pPr = p._element.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.SubElement(p._element, qn('a:pPr'))
    lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
    spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
    spcPct.set('val', str(int(line_spacing * 100000)))
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT1, thickness=3):
    """添加装饰性强调线"""
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = color
    line_shape.line.fill.background()
    return line_shape


def add_bullet_points(slide, left, top, width, height, points,
                      font_size=16, color=TEXT_BODY, spacing=Pt(10)):
    """添加要点列表（自定义符号）"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        clean_point = point.strip()
        # 保留已有编号
        if clean_point and clean_point[0].isdigit() and '.' in clean_point[:3]:
            p.text = clean_point
        else:
            symbol = BULLET_SYMBOLS[i % len(BULLET_SYMBOLS)]
            p.text = f"{symbol}  {clean_point}"

        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY_NAME
        p.font.bold = True
        p.space_after = spacing

        # 第一个字符（符号）用强调色
        if p.text and not p.text[0].isdigit():
            run = p.runs[0] if p.runs else None
            if run and len(run.text) > 0:
                # 重建 runs 以分色
                full_text = p.text
                p.clear()
                run1 = p.add_run()
                run1.text = full_text[0]
                run1.font.color.rgb = ACCENT1
                run1.font.size = Pt(font_size)
                run1.font.name = FONT_BODY_NAME
                run1.font.bold = True
                run2 = p.add_run()
                run2.text = full_text[1:]
                run2.font.color.rgb = color
                run2.font.size = Pt(font_size)
                run2.font.name = FONT_BODY_NAME
                run2.font.bold = True

    return txBox


def add_footer(slide, spec: SlideSpec, page_idx: int, total_pages: int):
    """添加底栏：页码 + 章节 + 进度条"""
    footer_y = Inches(7.05)
    footer_h = Inches(0.35)

    # 章节名
    section_short = spec.section.split('/')[0].strip() if '/' in spec.section else spec.section
    add_text_box(
        slide, Inches(0.5), footer_y, Inches(4), footer_h,
        section_short,
        font_size=9, color=TEXT_DIM, alignment=PP_ALIGN.LEFT
    )

    # 页码
    add_text_box(
        slide, Inches(5.5), footer_y, Inches(2.5), footer_h,
        f"{page_idx + 1} / {total_pages}",
        font_size=9, color=TEXT_DIM, alignment=PP_ALIGN.CENTER
    )

    # 累计时间
    if spec.cumulative:
        add_text_box(
            slide, Inches(9.5), footer_y, Inches(3.3), footer_h,
            spec.cumulative,
            font_size=9, color=TEXT_DIM, alignment=PP_ALIGN.RIGHT
        )

    # 进度条背景
    bar_y = Inches(7.38)
    bar_w = SLIDE_WIDTH
    bar_h = Pt(3)
    bar_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, bar_y, bar_w, bar_h)
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = PROGRESS_BG
    bar_bg.line.fill.background()

    # 进度条前景
    progress = (page_idx + 1) / total_pages
    bar_fg_w = int(SLIDE_WIDTH * progress)
    if bar_fg_w > 0:
        bar_fg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, bar_y, bar_fg_w, bar_h)
        bar_fg.fill.solid()
        bar_fg.fill.fore_color.rgb = ACCENT1
        bar_fg.line.fill.background()


def add_speaker_notes(slide, spec: SlideSpec):
    """添加演讲词和讲者提醒到备注"""
    notes = slide.notes_slide
    notes_tf = notes.notes_text_frame
    parts = []
    if spec.speech:
        parts.append(f"【演讲词】\n{spec.speech}")
    if spec.speaker_note:
        parts.append(f"【讲者提醒】\n{spec.speaker_note}")
    notes_tf.text = "\n\n".join(parts) if parts else ""


# ─── 幻灯片构建器 ─────────────────────────────────
def build_section_slide(prs, spec: SlideSpec, page_idx: int, total_pages: int):
    """构建幕标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_SECTION, RGBColor(0x0F, 0x12, 0x2A))

    # 获取章节编号
    section_num = ""
    section_color = ACCENT1
    for sec_name, sec_info in SECTIONS.items():
        if sec_name in spec.section or sec_name in spec.title:
            section_num = sec_info["num"]
            section_color = sec_info["color"]
            break

    # 章节编号装饰（大号、用较暗颜色模拟半透明效果）
    if section_num:
        # 用极暗版本的章节色模拟半透明
        dim_color = RGBColor(0x15, 0x2A, 0x25) if section_color == ACCENT1 else RGBColor(0x2A, 0x18, 0x10)
        add_text_box(
            slide, Inches(0.5), Inches(0.5), Inches(4), Inches(4),
            section_num,
            font_size=120, color=dim_color,
            bold=True, font_name=FONT_TITLE
        )

    # 主标题
    title_text = spec.ppt_title or spec.title
    add_text_box(
        slide, Inches(1.5), Inches(2.2), Inches(10), Inches(2),
        title_text,
        font_size=44, color=TEXT_WHITE, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE
    )

    # 标题下方强调线
    add_accent_line(slide, Inches(5.5), Inches(4.3), Inches(2.3), section_color, thickness=4)

    # 副标题
    if spec.ppt_subtitle_points:
        subtitle_text = "\n".join(spec.ppt_subtitle_points)
        add_text_box(
            slide, Inches(2), Inches(4.8), Inches(9), Inches(1.5),
            subtitle_text,
            font_size=18, color=TEXT_BODY,
            alignment=PP_ALIGN.CENTER, line_spacing=1.5
        )

    add_footer(slide, spec, page_idx, total_pages)
    add_speaker_notes(slide, spec)
    return slide


def build_quote_card(prs, spec: SlideSpec, page_idx: int, total_pages: int):
    """构建金句卡片页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, QUOTE_BG1, QUOTE_BG2, angle=2700000)  # 左上→右下

    # 提取金句
    quote_text = spec.image_value if spec.image_value else spec.ppt_title
    if not quote_text:
        quote_text = spec.ppt_title

    # 装饰大引号 — 左上
    add_text_box(
        slide, Inches(1.0), Inches(0.8), Inches(2), Inches(2),
        "\u201C",
        font_size=120, color=RGBColor(0x00, 0xD4, 0xAA),
        bold=True, font_name="Georgia"
    )
    # 设置引号透明度
    # (简化处理，直接用较暗的颜色模拟透明效果)

    # 装饰大引号 — 右下
    add_text_box(
        slide, Inches(10.3), Inches(4.5), Inches(2), Inches(2),
        "\u201D",
        font_size=120, color=RGBColor(0x00, 0x80, 0x66),
        bold=True, font_name="Georgia"
    )

    # 金句主体
    # 根据文本长度调整字号
    if len(quote_text) > 40:
        q_size = 28
    elif len(quote_text) > 25:
        q_size = 32
    else:
        q_size = 36

    add_text_box(
        slide, Inches(2), Inches(2.0), Inches(9.3), Inches(3.5),
        quote_text,
        font_size=q_size, color=TEXT_WHITE, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE,
        line_spacing=1.6
    )

    # 如果有额外要点
    if spec.ppt_subtitle_points:
        points_text = "  |  ".join(spec.ppt_subtitle_points[:3])
        add_text_box(
            slide, Inches(2), Inches(5.5), Inches(9), Inches(1),
            points_text,
            font_size=13, color=TEXT_DIM,
            alignment=PP_ALIGN.CENTER
        )

    # 底部细线装饰
    add_accent_line(slide, Inches(4), Inches(5.3), Inches(5.3), ACCENT1, thickness=2)

    add_footer(slide, spec, page_idx, total_pages)
    add_speaker_notes(slide, spec)
    return slide


def build_diagram_slide(prs, spec: SlideSpec, page_idx: int, total_pages: int):
    """构建图表页（插入渲染后的 PNG）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    # 标题
    add_text_box(
        slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
        spec.ppt_title or spec.title,
        font_size=30, color=ACCENT1, bold=True, font_name=FONT_TITLE
    )
    add_accent_line(slide, Inches(0.6), Inches(1.1), Inches(3), ACCENT1, thickness=3)

    # 查找对应的图表 PNG
    diagram_file = DIAGRAM_MAP.get(spec.page_num)
    diagram_path = DIAGRAMS_RENDERED / diagram_file if diagram_file else None

    if diagram_path and diagram_path.exists():
        # 插入图表 PNG
        try:
            img = Image.open(diagram_path)
            img_w, img_h = img.size
            aspect = img_w / img_h

            # 图表占据大部分页面
            max_w = Inches(11.5)
            max_h = Inches(5.3)

            if aspect > (max_w / max_h):
                w = max_w
                h = int(w / aspect)
            else:
                h = max_h
                w = int(h * aspect)

            # 居中放置
            left = (SLIDE_WIDTH - w) // 2
            top = Inches(1.4)
            slide.shapes.add_picture(str(diagram_path), left, top, w, h)
        except Exception as e:
            _add_placeholder_box(slide, f"图表加载失败: {e}")
    else:
        _add_placeholder_box(slide, f"[图表待渲染] {spec.image_value}")

    # 要点（如果有，放在图表下方简短展示）
    if spec.ppt_subtitle_points:
        points_text = "  ·  ".join(spec.ppt_subtitle_points[:4])
        add_text_box(
            slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
            points_text,
            font_size=11, color=TEXT_DIM,
            alignment=PP_ALIGN.CENTER
        )

    add_footer(slide, spec, page_idx, total_pages)
    add_speaker_notes(slide, spec)
    return slide


def _add_placeholder_box(slide, text):
    """添加占位框（仅在缺少资产时使用）"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3), Inches(2.5), Inches(7), Inches(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x18, 0x1E, 0x30)
    shape.line.color.rgb = ACCENT1
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_DIM
    p.alignment = PP_ALIGN.CENTER


def build_screenshot_slide(prs, spec: SlideSpec, page_idx: int, total_pages: int):
    """构建截图内容页（左文右图布局 + 圆角阴影）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    img_path = SCREENSHOTS / spec.image_value
    has_image = img_path.exists()

    if has_image:
        # ── 左文右图布局 ──
        # 标题（左侧）
        add_text_box(
            slide, Inches(0.6), Inches(0.3), Inches(5.8), Inches(0.9),
            spec.ppt_title or spec.title,
            font_size=28, color=ACCENT1, bold=True, font_name=FONT_TITLE
        )
        add_accent_line(slide, Inches(0.6), Inches(1.2), Inches(2.5), ACCENT1, thickness=3)

        # 要点（左侧）
        if spec.ppt_subtitle_points:
            add_bullet_points(
                slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(4.5),
                spec.ppt_subtitle_points,
                font_size=15, color=TEXT_BODY, spacing=Pt(12)
            )

        # 截图（右侧，带圆角阴影）
        try:
            processed_img = add_rounded_shadow(img_path, radius=16, shadow_offset=6, shadow_blur=12)
            pil_img = Image.open(processed_img)
            proc_w, proc_h = pil_img.size
            aspect = proc_w / proc_h

            max_w = Inches(6.2)
            max_h = Inches(5.5)

            if aspect > (max_w / max_h):
                w = max_w
                h = int(w / aspect)
            else:
                h = max_h
                w = int(h * aspect)

            left = SLIDE_WIDTH - w - Inches(0.4)
            top = Inches(1.0)

            processed_img.seek(0)
            slide.shapes.add_picture(processed_img, left, top, w, h)
        except Exception as e:
            add_text_box(
                slide, Inches(7), Inches(2.5), Inches(5.5), Inches(2),
                f"[图片: {img_path.name}]",
                font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER
            )
    else:
        # ── 全宽文字布局（无图片） ──
        add_text_box(
            slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(1),
            spec.ppt_title or spec.title,
            font_size=32, color=ACCENT1, bold=True, font_name=FONT_TITLE
        )
        add_accent_line(slide, Inches(0.8), Inches(1.35), Inches(3), ACCENT1, thickness=3)

        if spec.ppt_subtitle_points:
            add_bullet_points(
                slide, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                spec.ppt_subtitle_points,
                font_size=18, color=TEXT_BODY, spacing=Pt(14)
            )

        # 提示截图缺失
        add_text_box(
            slide, Inches(8), Inches(6.0), Inches(4.5), Inches(0.4),
            f"⚠ 截图缺失: {spec.image_value}",
            font_size=9, color=RGBColor(0xFF, 0x44, 0x44)
        )

    add_footer(slide, spec, page_idx, total_pages)
    add_speaker_notes(slide, spec)
    return slide


def build_dense_text_slide(prs, spec: SlideSpec, page_idx: int, total_pages: int):
    """构建高密度文字页：14pt字体，6-10条要点"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    # 标题（稍小以留更多空间给要点）
    add_text_box(
        slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
        spec.ppt_title or spec.title,
        font_size=26, color=ACCENT1, bold=True, font_name=FONT_TITLE
    )
    add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(2.5), ACCENT1, thickness=3)

    # 高密度要点
    if spec.ppt_subtitle_points:
        points = spec.ppt_subtitle_points
        n = len(points)
        # 单栏或双栏取决于要点数
        if n <= 6:
            add_bullet_points(
                slide, Inches(0.6), Inches(1.2), Inches(12), Inches(5.5),
                points,
                font_size=14, color=TEXT_BODY, spacing=Pt(8)
            )
        else:
            # 双栏布局
            mid = (n + 1) // 2
            add_bullet_points(
                slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(5.5),
                points[:mid],
                font_size=14, color=TEXT_BODY, spacing=Pt(8)
            )
            add_bullet_points(
                slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.5),
                points[mid:],
                font_size=14, color=TEXT_BODY, spacing=Pt(8)
            )

    add_footer(slide, spec, page_idx, total_pages)
    add_speaker_notes(slide, spec)
    return slide


def build_table_slide(prs, spec: SlideSpec, page_idx: int, total_pages: int):
    """构建表格对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    # 标题
    add_text_box(
        slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
        spec.ppt_title or spec.title,
        font_size=26, color=ACCENT1, bold=True, font_name=FONT_TITLE
    )
    add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(2.5), ACCENT1, thickness=3)

    # 表格
    table_rows = spec.table_data
    if not table_rows:
        # 如果没有解析到表格数据，降级为dense
        if spec.ppt_subtitle_points:
            add_bullet_points(
                slide, Inches(0.6), Inches(1.2), Inches(12), Inches(5.5),
                spec.ppt_subtitle_points,
                font_size=14, color=TEXT_BODY, spacing=Pt(8)
            )
    else:
        rows = len(table_rows)
        cols = max(len(r) for r in table_rows) if table_rows else 2
        # 确保所有行等宽
        for r in table_rows:
            while len(r) < cols:
                r.append("")

        table_left = Inches(0.6)
        table_top = Inches(1.3)
        table_width = Inches(12)
        table_height = Inches(5.2)

        table_shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
        table = table_shape.table

        # 设置列宽均分
        col_width = int(table_width / cols)
        for ci in range(cols):
            table.columns[ci].width = col_width

        for ri, row_data in enumerate(table_rows):
            for ci, cell_text in enumerate(row_data):
                cell = table.cell(ri, ci)
                cell.text = cell_text
                # 样式
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(13 if ri > 0 else 14)
                    paragraph.font.name = FONT_BODY_NAME
                    paragraph.font.color.rgb = TEXT_WHITE if ri == 0 else TEXT_BODY
                    paragraph.font.bold = (ri == 0)
                # 背景色
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
                srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
                if ri == 0:
                    srgb.set('val', '0A1628')  # 深色表头
                elif ri % 2 == 0:
                    srgb.set('val', '111B2E')  # 偶数行
                else:
                    srgb.set('val', '0E1626')  # 奇数行

    # 补充要点（表格下方）
    if spec.ppt_subtitle_points and spec.table_data:
        points_text = "  ·  ".join(spec.ppt_subtitle_points[:3])
        add_text_box(
            slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
            points_text,
            font_size=11, color=TEXT_DIM,
            alignment=PP_ALIGN.CENTER
        )

    add_footer(slide, spec, page_idx, total_pages)
    add_speaker_notes(slide, spec)
    return slide


def build_text_slide(prs, spec: SlideSpec, page_idx: int, total_pages: int):
    """构建纯文字内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BG_PRIMARY, BG_SECONDARY)

    # 标题
    add_text_box(
        slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(1),
        spec.ppt_title or spec.title,
        font_size=32, color=ACCENT1, bold=True, font_name=FONT_TITLE
    )
    add_accent_line(slide, Inches(0.8), Inches(1.35), Inches(3), ACCENT1, thickness=3)

    # 要点
    if spec.ppt_subtitle_points:
        add_bullet_points(
            slide, Inches(0.8), Inches(1.8), Inches(11), Inches(4.8),
            spec.ppt_subtitle_points,
            font_size=20, color=TEXT_BODY, spacing=Pt(16)
        )

    add_footer(slide, spec, page_idx, total_pages)
    add_speaker_notes(slide, spec)
    return slide


# ─── 页面类型判断 ─────────────────────────────────
def is_section_title(spec: SlideSpec) -> bool:
    """判断是否为幕标题页"""
    title = spec.title
    # 包含 "第X幕" 或 "序幕" 的标题页
    is_act = bool(re.search(r'第[一二三四五六七]幕', title) or '序幕' in title)
    # diagram-placeholder 页应优先渲染图表，不算"无视觉"
    if spec.image_type == 'diagram-placeholder':
        return False
    has_no_visual = spec.image_type in ('text-only', '')
    is_title_page = '标题页' in title or '开场' in title
    return is_act and (has_no_visual or is_title_page)


# ─── 主构建流程 ───────────────────────────────────
def build_presentation(all_slides: list[SlideSpec]) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    total = len(all_slides)
    stats = {"section": 0, "quote": 0, "diagram": 0, "screenshot": 0, "text": 0, "dense": 0, "table": 0}

    for i, spec in enumerate(all_slides):
        slide_type = ""

        # Layout字段优先
        if spec.layout == "dense":
            build_dense_text_slide(prs, spec, i, total)
            slide_type = "dense"
            stats["dense"] += 1
        elif spec.layout == "table":
            build_table_slide(prs, spec, i, total)
            slide_type = "table"
            stats["table"] += 1
        elif spec.layout == "section" or is_section_title(spec):
            build_section_slide(prs, spec, i, total)
            slide_type = "section"
            stats["section"] += 1
        elif spec.image_type == "quote-card":
            build_quote_card(prs, spec, i, total)
            slide_type = "quote"
            stats["quote"] += 1
        elif spec.image_type == "diagram-placeholder":
            build_diagram_slide(prs, spec, i, total)
            slide_type = "diagram"
            stats["diagram"] += 1
        elif spec.image_type == "screenshot":
            build_screenshot_slide(prs, spec, i, total)
            slide_type = "screenshot"
            stats["screenshot"] += 1
        else:
            build_text_slide(prs, spec, i, total)
            slide_type = "text"
            stats["text"] += 1

        print(f"  [{i+1:3d}/{total}] 页{spec.page_num:>4s} [{slide_type:10s}] {spec.ppt_title[:35] if spec.ppt_title else spec.title[:35]}")

    return prs, stats


def main():
    print("=" * 65)
    print("  PPTX Build v2 — AI 科普演讲（高质量重建）")
    print("=" * 65)

    # 1. 解析
    all_slides = []
    for spec_file in SPEC_FILES:
        print(f"\n  解析: {spec_file.name}")
        slides = parse_spec_file(spec_file)
        print(f"    → {len(slides)} 页")
        all_slides.extend(slides)

    print(f"\n  总计: {len(all_slides)} 页")

    # 2. 类型统计
    type_counts = {}
    for s in all_slides:
        t = s.image_type or "unknown"
        type_counts[t] = type_counts.get(t, 0) + 1
    print("\n  配图类型统计:")
    for t, c in sorted(type_counts.items()):
        print(f"    {t}: {c}")

    # 3. 检查截图
    missing = []
    for s in all_slides:
        if s.image_type == "screenshot":
            path = SCREENSHOTS / s.image_value
            if not path.exists():
                missing.append((s.page_num, s.image_value))
    if missing:
        print(f"\n  ⚠ {len(missing)} 张截图缺失:")
        for pn, fn in missing[:5]:
            print(f"    页 {pn}: {fn}")
        if len(missing) > 5:
            print(f"    ... 还有 {len(missing)-5} 张")

    # 4. 检查图表
    diagram_status = []
    for page, fname in DIAGRAM_MAP.items():
        path = DIAGRAMS_RENDERED / fname
        status = "✓" if path.exists() else "✗"
        diagram_status.append((page, fname, status))
    print(f"\n  图表状态:")
    for page, fname, status in diagram_status:
        print(f"    {status} 页{page}: {fname}")

    # 5. 构建
    print("\n  开始构建 PPTX...")
    print("-" * 65)
    prs, stats = build_presentation(all_slides)
    print("-" * 65)

    # 6. 保存
    prs.save(str(OUTPUT_PPTX))
    size_mb = OUTPUT_PPTX.stat().st_size / (1024 * 1024)

    print(f"\n  ✅ 已保存: {OUTPUT_PPTX}")
    print(f"     文件大小: {size_mb:.1f} MB")
    print(f"     总页数: {len(all_slides)}")
    print(f"\n  构建统计:")
    for k, v in stats.items():
        print(f"    {k}: {v}")

    # 7. 质量检查
    print(f"\n  质量检查:")
    placeholder_count = sum(1 for s in all_slides
                           if s.image_type == "diagram-placeholder"
                           and not (DIAGRAMS_RENDERED / DIAGRAM_MAP.get(s.page_num, "")).exists())
    print(f"    未替换的 diagram-placeholder: {placeholder_count}")
    print(f"    截图缺失: {len(missing)}")
    print(f"    文件 < 30MB: {'✓' if size_mb < 30 else '✗'} ({size_mb:.1f} MB)")
    print("=" * 65)


if __name__ == "__main__":
    main()
