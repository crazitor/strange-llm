#!/usr/bin/env python3
"""
AI科普演讲 PPTX 生成器 v2 — 需求驱动版
支持5种视觉类型: screenshot / text-only / section-cover / quote-card / diagram
"""

import re
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 路径 ──────────────────────────────────────────
BASE_DIR = Path(__file__).resolve().parent.parent.parent
OUTPUT_DIR = BASE_DIR / "output"
SCREENSHOTS_DIR = OUTPUT_DIR / "screenshots"
SPEC_PART1 = OUTPUT_DIR / "slides_spec_part1.md"
SPEC_PART2 = OUTPUT_DIR / "slides_spec_part2.md"
OUTPUT_PPTX = OUTPUT_DIR / "presentation.pptx"

# ── 颜色 ──────────────────────────────────────────
BG = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
DIM_GRAY = RGBColor(0x55, 0x55, 0x55)
TEAL = RGBColor(0x4E, 0xC9, 0xB0)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
WARM = RGBColor(0xFF, 0xA0, 0x60)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Microsoft YaHei"

# ── 需求驱动映射 ──────────────────────────────────
# 截图映射: page_num -> 截图文件名
SCREENSHOT_MAP = {
    3:  "video2_t0029_输出博士也极其呆板.jpg",
    4:  "video2_t0130_然后我又按照重要程度分成了里程碑.jpg",
    6:  "video2_t1030_交互协议是否规范.jpg",
    7:  "video2_t0059_今天这期视频我们就来一起探讨下这个问题.jpg",
    8:  "video1_t0100_那如果把角色区分一下.jpg",
    9:  "video1_t0050_恭喜你发明了今天的第一个新词儿.jpg",
    10: "video1_t0140_有的部分是最终的指示.jpg",
    11: "video1_t0115_只能一问一答.jpg",
    12: "video1_t0159_作为上下文信息.jpg",
    13: "video1_t0205_给这些特殊的善寒文信息起了个新词叫memory.jpg",
    15: "video1_t0229_你发现的第一个问题就是.jpg",
    16: "video1_t0254_但很快你就发现这样好像显得自己有点蠢.jpg",
    17: "video1_t0330_那从现在的视角回看.jpg",
    19: "video1_t0435_我们聚焦于agent和大模型的对话过程来看.jpg",
    20: "video1_t0345_把语义相近的片段找出来.jpg",
    21: "video1_t0409_连RG也算是search的一种了.jpg",
    22: "video1_t0459_其实呢就是个约定罢了.jpg",
    23: "video1_t0529_翻译过来叫模型上下文协议.jpg",
    24: "video1_t0524_也就是一套约定而已.jpg",
    25: "video1_t0549_同时呢别忘了还要给你这个老板传话.jpg",
    26: "video1_t0705_中间的翻译直接和大模型沟通即可.jpg",
    27: "video1_t0755_这个时候就又不好用程序来判断分支了.jpg",
    28: "video1_t0819_然后再按照要求完成任务.jpg",
    29: "video1_t0909_隔离子agent的产生的上下文不会保留在主A的中.jpg",
    30: "video1_t0930_就是有点拉垮.jpg",
    36: "video4_t0024_比如最著名的开源操作系统LINUX.jpg",
    37: "video4_t0319_最多可能在论文中描述一下思路.jpg",
    38: "video4_t0329_说不定这个能力就超越了所有大模型呢.jpg",
    39: "video4_t0345_那如果套盒指的是抄袭的权重.jpg",
    40: "video5_t0200_对不对.jpg",
    41: "video5_t0431_那这个时候模型的数量就只剩下八个了.jpg",
    44: "video2_t0130_然后我又按照重要程度分成了里程碑.jpg",
    45: "video2_t0159_已经为机器是否具备智能.jpg",
    46: "video2_t0300_我到现在还能清楚的记得.jpg",
    47: "video2_t0400_说实话.jpg",
    48: "video2_t0530_这期视频同时呢在刚刚的23年提到了.jpg",
    50: "video2_t0701_我们的主观世界和客观世界的节奏错位了.jpg",
    51: "video2_t0829_却可能数不清六根手指.jpg",
    # Part 2
    57: "video3_t0030_我的组合是腾讯云轻量服务器加openRT.jpg",
    58: "video3_t0100_嘻嘻嘻.jpg",
    59: "video3_t0230_因为opencloud有一大堆系统提示词.jpg",
    60: "video3_t0255_把今日的AI新闻整理成PDF发给我.jpg",
    61: "video3_t0359_实际上在后台其实就是增加了个chrome的配置文件.jpg",
    68: "video6_t0000_你觉得这四个AI生成的视频哪个难度最大.jpg",
    69: "video6_t0115_那就变成了我们熟悉的向量.jpg",
    70: "video6_t0139_这里面呢就有一些点是看起来合理的视频.jpg",
    71: "video6_t0229_但是呢这样做太难了.jpg",
    72: "video6_t0300_没错.jpg",
    73: "video6_t0320_把图片连带着时间簿以及对应的文字.jpg",
    74: "video6_t0345_这就是当今文生图和文生视频的底层架构.jpg",
    75: "video6_t0434_比如最新的cls20中.jpg",
    77: "video6_t0409_沿着某个方向移动到合理图片上的点的过程.jpg",
    78: "video6_t0430_现在的AI视频模型已经有相当多的.jpg",
    79: "video2_t0829_却可能数不清六根手指.jpg",
    80: "video2_t0900_也可能直接失效.jpg",
    81: "video2_t0930_这五类问题不可能被彻底消除.jpg",
}

# 章节封面页
SECTION_COVERS = {1, 18, 34, 43, 56, 67, 76, 85, 94, 101}

# 纯文字页
TEXT_ONLY = {2, 32, 35, 49, 64, 65, 83, 86, 88, 89, 90, 91, 92, 96, 97, 98, 100, 102}

# 金句卡片
QUOTE_CARDS = {
    5:  '全球震撼 → 集体耸肩\n\n"就这？"\n"感觉跟上一代也没啥区别"\n"说好的AGI呢？"',
    31: 'Token一定会越来越便宜\n\n光谱上的平衡点不断向"柔性"移动\nWorkflow（写死） → Agent（灵活）',
    42: '模型结构 = 图纸\n权重 = 房子\n训练 = 施工过程\n\n"战略上要藐视技术，战术上要重视技术"',
    63: '"虽然他的每一步思考都非常合理，\n都是遇到问题然后解决问题，\n但现实生活中如果真有人干这种事，\n那你多半觉得他是个疯子。"\n\n—— 而且这个人还要按照付出多少努力来收费',
    82: '"越来越难以发现他胡说八道的地方了，\n也就是他胡说八道的能力也越来越强了。"',
    93: '"你应该非常努力地学习用AI，\n然后非常努力地去做那些\nAI无法完成的内容。"',
    103: '"意义从来不在可能性的空间里，\n而在于选择本身。"\n\n可能性空间：数学的、冰冷的、无限的\n选择：人的、滚烫的、有限的',
    104: '"人之所以为人，\n不是因为我们能够创造出\n宇宙中从未出现过的组合，\n而是因为我们会在无穷多种可能性中\n固执地选择其中那一个。"\n\n有限性就是人类独有的浪漫。',
    105: '你的选择，AI做不了。\n只有你能做。',
}

# diagram 页（标记即可，这些页用text+简单shape构建）
DIAGRAM_PAGES = {14, 33, 52, 62, 66, 84, 87, 95, 99}


# ── 解析逻辑 ─────────────────────────────────────
def parse_slides(filepath: Path) -> list[dict]:
    text = filepath.read_text(encoding="utf-8")
    page_re = re.compile(r"^### 页 (\d+): (.+)$", re.MULTILINE)
    matches = list(page_re.finditer(text))
    pages = []

    for i, m in enumerate(matches):
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        block = text[start:end].strip()

        slide = {
            "page_num": int(m.group(1)),
            "page_title": m.group(2).strip(),
            "chapter": "",
            "time": "",
            "cumulative_time": "",
            "title": "",
            "points": [],
            "speaker_script": "",
            "speaker_reminder": "",
        }

        ch = re.search(r"\*\*章节\*\*:\s*(.+)", block)
        if ch: slide["chapter"] = ch.group(1).strip()

        tm = re.search(r"\*\*时间\*\*:\s*(.+)", block)
        if tm: slide["time"] = tm.group(1).strip()

        ct = re.search(r"\*\*累计时间\*\*:\s*(.+)", block)
        if ct: slide["cumulative_time"] = ct.group(1).strip()

        ppt_block = re.search(
            r"\*\*PPT显示内容\*\*:\s*\n(.*?)(?=\n\*\*演讲词\*\*)", block, re.DOTALL
        )
        if ppt_block:
            ppt_text = ppt_block.group(1)
            title_m = re.search(r"- 标题:\s*(.+)", ppt_text)
            if title_m: slide["title"] = title_m.group(1).strip()
            slide["points"] = re.findall(r"^\s{2,}-\s+(.+)$", ppt_text, re.MULTILINE)
            slide["points"] = [p.strip() for p in slide["points"]]

        script_m = re.search(
            r"\*\*演讲词\*\*:\s*\n>\s*(.+?)(?=\n\n\*\*讲者提醒|\n---|\Z)", block, re.DOTALL
        )
        if script_m:
            raw = script_m.group(1).strip()
            lines = [l.strip().lstrip("> ").strip() for l in raw.split("\n")]
            slide["speaker_script"] = " ".join(lines)

        rem_m = re.search(r"\*\*讲者提醒\*\*:\s*(.+?)(?=\n---|\Z)", block, re.DOTALL)
        if rem_m: slide["speaker_reminder"] = rem_m.group(1).strip()

        pages.append(slide)
    return pages


def get_visual_type(page_num):
    if page_num in SECTION_COVERS: return "section-cover"
    if page_num in QUOTE_CARDS: return "quote-card"
    if page_num in DIAGRAM_PAGES: return "diagram"
    if page_num in SCREENSHOT_MAP: return "screenshot"
    if page_num in TEXT_ONLY: return "text-only"
    return "text-only"  # 默认


# ── 渲染函数 ─────────────────────────────────────
def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def text_box(slide, left, top, w, h, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    return tb


def multi_text(slide, left, top, w, h, lines, size=18, color=LIGHT_GRAY,
               line_spacing=Pt(8)):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        is_quote = line.startswith('"') or line.startswith('"') or line.startswith("「")
        p.text = f"• {line}" if not is_quote else line
        p.font.size = Pt(size)
        p.font.color.rgb = TEAL if is_quote else color
        p.font.italic = is_quote
        p.font.name = FONT
        p.space_after = line_spacing
    return tb


def add_notes(slide, sd):
    parts = []
    if sd["time"]: parts.append(f"⏱ {sd['time']}")
    if sd["cumulative_time"]: parts.append(f"📍 {sd['cumulative_time']}")
    if sd["speaker_script"]:
        parts.append("")
        parts.append(sd["speaker_script"])
    if sd["speaker_reminder"]:
        parts.append(f"\n⚠️ {sd['speaker_reminder']}")
    slide.notes_slide.notes_text_frame.text = "\n".join(parts)


def page_label(slide, sd):
    if sd["chapter"]:
        text_box(slide, Inches(0.6), Inches(0.15), Inches(5), Inches(0.3),
                 sd["chapter"], size=10, color=DIM_GRAY)
    text_box(slide, Inches(12.2), Inches(7), Inches(0.8), Inches(0.3),
             str(sd["page_num"]), size=10, color=DIM_GRAY, align=PP_ALIGN.RIGHT)


# ── 5种幻灯片类型 ────────────────────────────────

def render_section_cover(prs, sd):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title = sd["title"] or sd["page_title"]
    text_box(s, Inches(1), Inches(2.2), Inches(11.333), Inches(1.5),
             title, size=48, bold=True, align=PP_ALIGN.CENTER)
    if sd["points"]:
        text_box(s, Inches(2), Inches(4.2), Inches(9.333), Inches(1.5),
                 "\n".join(sd["points"]), size=22, color=LIGHT_GRAY,
                 align=PP_ALIGN.CENTER)
    page_label(s, sd)
    add_notes(s, sd)


def render_text_only(prs, sd):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title = sd["title"] or sd["page_title"]
    text_box(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
             title, size=32, bold=True)
    if sd["points"]:
        multi_text(s, Inches(1), Inches(1.8), Inches(11), Inches(5),
                   sd["points"], size=20)
    page_label(s, sd)
    add_notes(s, sd)


def render_screenshot(prs, sd):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    pn = sd["page_num"]
    img_file = SCREENSHOT_MAP.get(pn, "")
    img_path = SCREENSHOTS_DIR / img_file

    title = sd["title"] or sd["page_title"]

    if img_path.exists():
        # 左文右图布局
        text_box(s, Inches(0.8), Inches(0.5), Inches(7), Inches(1),
                 title, size=30, bold=True)
        if sd["points"]:
            multi_text(s, Inches(1), Inches(1.7), Inches(6.5), Inches(5),
                       sd["points"], size=17)
        try:
            s.shapes.add_picture(str(img_path),
                                 Inches(8.2), Inches(1.2),
                                 Inches(4.6), Inches(5.5))
        except Exception as e:
            text_box(s, Inches(8.2), Inches(3), Inches(4.6), Inches(1),
                     f"[{img_file}]", size=11, color=DIM_GRAY,
                     align=PP_ALIGN.CENTER)
    else:
        # 降级为 text-only
        text_box(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
                 title, size=32, bold=True)
        if sd["points"]:
            multi_text(s, Inches(1), Inches(1.8), Inches(11), Inches(5),
                       sd["points"], size=20)

    page_label(s, sd)
    add_notes(s, sd)


def render_quote_card(prs, sd):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    pn = sd["page_num"]
    quote = QUOTE_CARDS.get(pn, "")
    title = sd["title"] or sd["page_title"]

    # 小标题
    text_box(s, Inches(1), Inches(0.8), Inches(11.333), Inches(0.5),
             title, size=18, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # 金句居中
    tb = s.shapes.add_textbox(Inches(1.5), Inches(1.8),
                              Inches(10.333), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(quote.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
        if line.startswith('"') or line.startswith('"'):
            p.font.size = Pt(28)
            p.font.color.rgb = GOLD
            p.font.bold = True
        elif line == "":
            p.font.size = Pt(12)
        else:
            p.font.size = Pt(22)
            p.font.color.rgb = LIGHT_GRAY
        p.space_after = Pt(6)

    page_label(s, sd)
    add_notes(s, sd)


def render_diagram(prs, sd):
    """Diagram 页面: 用标题+要点+简单shape占位"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title = sd["title"] or sd["page_title"]

    text_box(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
             title, size=32, bold=True)

    pn = sd["page_num"]

    # 根据页面内容构建简化diagram
    if pn == 14:  # 四个概念递进
        labels = ["LLM\n文字接龙机器", "Prompt\n你的指令", "Context\n带上聊天记录", "Memory\n每次都贴的规矩"]
        for i, lab in enumerate(labels):
            left = Inches(0.8 + i * 3.1)
            shape = s.shapes.add_shape(1, left, Inches(2.5), Inches(2.6), Inches(2))  # 1=rectangle
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x4E)
            shape.line.color.rgb = TEAL
            shape.line.width = Pt(2)
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = lab
            p.font.size = Pt(16)
            p.font.color.rgb = WHITE
            p.font.name = FONT
            p.alignment = PP_ALIGN.CENTER
        text_box(s, Inches(1), Inches(5.5), Inches(11), Inches(1),
                 '套路: 遇到问题 → 想个办法 → 起个洋气名字', size=18, color=TEAL,
                 align=PP_ALIGN.CENTER)

    elif pn == 33:  # 小L进化路径
        steps = ["文字接龙\n呆瓜", "+Prompt\n能听指令", "+Context\n能对话", "+Memory\n有规矩", "+Agent\n能做事", "+Skill\n管家"]
        for i, step in enumerate(steps):
            left = Inches(0.4 + i * 2.1)
            shape = s.shapes.add_shape(1, left, Inches(2.5), Inches(1.8), Inches(2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x4E)
            shape.line.color.rgb = TEAL
            shape.line.width = Pt(2)
            tf = shape.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = step; p.font.size = Pt(14)
            p.font.color.rgb = WHITE; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
        text_box(s, Inches(1), Inches(5.5), Inches(11), Inches(1),
                 '内核从未改变: 文字接龙机器', size=20, color=WARM, align=PP_ALIGN.CENTER)

    elif pn == 52:  # 前半场回顾
        cards = [("底裤", "文字接龙=LLM本质"), ("脚手架", "Agent+概念体系"),
                 ("图纸vs施工", "开源真相/万物归一"), ("时间线", "从惊艳到质变")]
        for i, (t, sub) in enumerate(cards):
            left = Inches(0.6 + i * 3.15)
            shape = s.shapes.add_shape(1, left, Inches(2), Inches(2.8), Inches(3))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x45)
            shape.line.color.rgb = TEAL; shape.line.width = Pt(2)
            tf = shape.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = t; p.font.size = Pt(22); p.font.bold = True
            p.font.color.rgb = TEAL; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
            p2 = tf.add_paragraph(); p2.text = sub; p2.font.size = Pt(14)
            p2.font.color.rgb = LIGHT_GRAY; p2.font.name = FONT; p2.alignment = PP_ALIGN.CENTER
        text_box(s, Inches(1), Inches(6), Inches(11), Inches(0.8),
                 '中场休息 10 分钟后，我们继续', size=22, color=LIGHT_GRAY,
                 align=PP_ALIGN.CENTER)

    elif pn == 62:  # 买充电器流程
        steps = ["老板问\n几号?", "掏手机\n没电", "买充电器\n超市关门", "打车远\n没法付款", "去银行\n余额不够", "抵押\n手机..."]
        for i, step in enumerate(steps):
            left = Inches(0.4 + i * 2.1)
            shape = s.shapes.add_shape(1, left, Inches(2.5), Inches(1.8), Inches(2.2))
            shape.fill.solid()
            # 颜色从绿渐变到红
            r = min(255, 80 + i * 35); g = max(50, 200 - i * 30)
            shape.fill.fore_color.rgb = RGBColor(r, g, 0x30)
            shape.line.color.rgb = WHITE; shape.line.width = Pt(1)
            tf = shape.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = step; p.font.size = Pt(14)
            p.font.color.rgb = WHITE; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
        text_box(s, Inches(1), Inches(5.5), Inches(11), Inches(1),
                 '每一步都很合理，但整体是个疯子', size=22, color=WARM, align=PP_ALIGN.CENTER)

    elif pn == 66:  # 交通灯矩阵
        rows = [("简单任务", "完美但昂贵", RGBColor(0x40, 0xC0, 0x60)),
                ("复杂任务", "能做但粗糙", RGBColor(0xE0, 0xC0, 0x40)),
                ("自主任务", "执着但失控", RGBColor(0xE0, 0x50, 0x50))]
        for i, (task, result, clr) in enumerate(rows):
            top = Inches(2 + i * 1.5)
            shape = s.shapes.add_shape(1, Inches(2), top, Inches(4), Inches(1.2))
            shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x45)
            shape.line.color.rgb = clr; shape.line.width = Pt(2)
            tf = shape.text_frame; p = tf.paragraphs[0]; p.text = task
            p.font.size = Pt(20); p.font.color.rgb = WHITE; p.font.name = FONT
            shape2 = s.shapes.add_shape(1, Inches(7), top, Inches(4), Inches(1.2))
            shape2.fill.solid(); shape2.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x45)
            shape2.line.color.rgb = clr; shape2.line.width = Pt(2)
            tf2 = shape2.text_frame; p2 = tf2.paragraphs[0]; p2.text = result
            p2.font.size = Pt(20); p2.font.color.rgb = clr; p2.font.name = FONT
        text_box(s, Inches(1), Inches(6.2), Inches(11), Inches(0.8),
                 'Agent是2025年AI落地的最大战场', size=18, color=TEAL, align=PP_ALIGN.CENTER)

    elif pn == 84:  # 三道枷锁
        circles = [("认知层", "解空间诅咒+OOD", Inches(4.5), Inches(2.2), Inches(4.3), Inches(3.5)),
                   ("架构层", "五类不可消除缺陷", Inches(3.5), Inches(1.7), Inches(6.3), Inches(4.5)),
                   ("社会层", "生态壁垒+利益博弈", Inches(2.5), Inches(1.2), Inches(8.3), Inches(5.5))]
        for label, sub, left, top, w, h in reversed(circles):
            shape = s.shapes.add_shape(1, left, top, w, h)
            shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x42)
            shape.line.color.rgb = TEAL; shape.line.width = Pt(2)
        # 文字标签覆盖
        for i, (label, sub, left, top, w, h) in enumerate(circles):
            y = Inches(2.2 + i * 1.3)
            text_box(s, Inches(1), y, Inches(4), Inches(0.5),
                     f"{label}: {sub}", size=16, color=WHITE if i == 0 else LIGHT_GRAY)
        text_box(s, Inches(1), Inches(6), Inches(11), Inches(0.8),
                 '突破需要: 新范式? 新架构? 新生态?', size=18, color=WARM, align=PP_ALIGN.CENTER)

    elif pn == 87:  # Skill vs MCP
        for i, (title_t, sub, items) in enumerate([
            ("Skill = 菜谱", "脑子里的知识", ["提示词指导LLM做什么", "存在本地配置目录", "知识层"]),
            ("MCP = 厨具", "手里的工具", ["接入外部工具的协议", "标准化接口规范", "协议层"])
        ]):
            left = Inches(1 + i * 6)
            shape = s.shapes.add_shape(1, left, Inches(2), Inches(5), Inches(4))
            shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x42)
            shape.line.color.rgb = TEAL; shape.line.width = Pt(2)
            tf = shape.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = title_t; p.font.size = Pt(24); p.font.bold = True
            p.font.color.rgb = TEAL; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
            p2 = tf.add_paragraph(); p2.text = sub; p2.font.size = Pt(16)
            p2.font.color.rgb = LIGHT_GRAY; p2.font.name = FONT; p2.alignment = PP_ALIGN.CENTER
            for item in items:
                p3 = tf.add_paragraph(); p3.text = f"• {item}"; p3.font.size = Pt(14)
                p3.font.color.rgb = LIGHT_GRAY; p3.font.name = FONT
        text_box(s, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 '完全不同层次，互不替代', size=18, color=WARM, align=PP_ALIGN.CENTER)

    elif pn == 95:  # Agent vs SOAR
        for i, (col_title, steps) in enumerate([
            ("Agent 流程", ["理解意图", "选择工具", "执行", "反馈"]),
            ("SOAR 流程", ["告警输入", "剧本匹配", "工具调用", "响应处置"])
        ]):
            left = Inches(1.5 + i * 5.5)
            text_box(s, left, Inches(1.8), Inches(4), Inches(0.6),
                     col_title, size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
            for j, step in enumerate(steps):
                top = Inches(2.6 + j * 1)
                shape = s.shapes.add_shape(1, left + Inches(0.5), top, Inches(3), Inches(0.7))
                shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x45)
                shape.line.color.rgb = TEAL; shape.line.width = Pt(1)
                tf = shape.text_frame; p = tf.paragraphs[0]; p.text = step
                p.font.size = Pt(16); p.font.color.rgb = WHITE; p.font.name = FONT
                p.alignment = PP_ALIGN.CENTER
        text_box(s, Inches(1), Inches(6.2), Inches(11), Inches(0.6),
                 '底层逻辑完全一致', size=20, color=GOLD, align=PP_ALIGN.CENTER)

    elif pn == 99:  # 五个赋能方向
        dirs = [("漏洞扫描", "AI辅助漏洞分析+PoC生成"),
                ("威胁情报", "AI驱动关联+趋势预测"),
                ("安全运营", "智能告警降噪+自动研判"),
                ("安全培训", "AI生成攻防演练场景"),
                ("客户服务", "AI辅助安全咨询+方案推荐")]
        for i, (d, desc) in enumerate(dirs):
            top = Inches(1.8 + i * 1.05)
            shape = s.shapes.add_shape(1, Inches(1.5), top, Inches(3), Inches(0.8))
            shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x45)
            shape.line.color.rgb = TEAL; shape.line.width = Pt(2)
            tf = shape.text_frame; p = tf.paragraphs[0]; p.text = d
            p.font.size = Pt(18); p.font.color.rgb = TEAL; p.font.bold = True; p.font.name = FONT
            text_box(s, Inches(5), top + Inches(0.1), Inches(7), Inches(0.6),
                     desc, size=16, color=LIGHT_GRAY)
    else:
        # 通用 diagram 降级
        if sd["points"]:
            multi_text(s, Inches(1), Inches(1.8), Inches(11), Inches(5),
                       sd["points"], size=18)

    page_label(s, sd)
    add_notes(s, sd)


# ── 中场休息过渡页 ───────────────────────────────
def intermission(prs):
    # 页53
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    text_box(s, Inches(1), Inches(2.5), Inches(11.333), Inches(1.5),
             "中场休息", size=54, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, Inches(2), Inches(4.5), Inches(9.333), Inches(1),
             "休息 10 分钟，稍后继续", size=28, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    text_box(s, Inches(12.2), Inches(7), Inches(0.8), Inches(0.3),
             "53", size=10, color=DIM_GRAY, align=PP_ALIGN.RIGHT)
    s.notes_slide.notes_text_frame.text = "⏱ 10:00 中场休息"

    # 页54 关键词
    s2 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s2)
    text_box(s2, Inches(1), Inches(1), Inches(11.333), Inches(0.8),
             "前半场关键词", size=36, bold=True, align=PP_ALIGN.CENTER)
    kw = ("LLM  •  Token  •  Transformer  •  Agent  •  MCP\n"
          "Function Calling  •  RAG  •  Embedding  •  开源\n"
          "Scaling Law  •  蒸馏  •  GPT-4  •  DeepSeek\n"
          "多模态  •  推理模型  •  Context Window")
    text_box(s2, Inches(1.5), Inches(2.5), Inches(10.333), Inches(4),
             kw, size=28, color=TEAL, align=PP_ALIGN.CENTER)
    text_box(s2, Inches(12.2), Inches(7), Inches(0.8), Inches(0.3),
             "54", size=10, color=DIM_GRAY, align=PP_ALIGN.RIGHT)

    # 页55 下半场预告
    s3 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s3)
    text_box(s3, Inches(1), Inches(2), Inches(11.333), Inches(1),
             "下半场预告", size=40, bold=True, align=PP_ALIGN.CENTER)
    preview = ("• Agent 实战: 从查内存到买充电器的疯子\n"
               "• AI 视频生成的底层逻辑与诅咒\n"
               "• 能力天花板: 五类不可消除的缺陷\n"
               "• 纠偏专题: 八个深度思考\n"
               "• 安全业务关联: SOAR / 漏扫 / 自动渗透\n"
               "• 终章: 有限性的浪漫")
    text_box(s3, Inches(2), Inches(3.8), Inches(9.333), Inches(3),
             preview, size=22, color=LIGHT_GRAY)
    text_box(s3, Inches(12.2), Inches(7), Inches(0.8), Inches(0.3),
             "55", size=10, color=DIM_GRAY, align=PP_ALIGN.RIGHT)
    s3.notes_slide.notes_text_frame.text = (
        "【演讲词】好，欢迎回来。下半场会更精彩——"
        "Agent实战、AI视频底层逻辑、能力天花板、八个纠偏、"
        "安全业务关联、以及关于'人'的终章。准备好了吗？我们继续。")


# ── 主流程 ────────────────────────────────────────
def main():
    print("=" * 60)
    print("AI科普演讲 PPTX v2 — 需求驱动版")
    print("=" * 60)

    slides_p1 = parse_slides(SPEC_PART1)
    slides_p2 = parse_slides(SPEC_PART2)
    print(f"解析: Part1 {len(slides_p1)} 页, Part2 {len(slides_p2)} 页")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    stats = {"screenshot": 0, "text-only": 0, "section-cover": 0,
             "quote-card": 0, "diagram": 0}
    img_ok = 0; img_miss = 0

    for sd in slides_p1:
        vtype = get_visual_type(sd["page_num"])
        stats[vtype] += 1
        if vtype == "section-cover": render_section_cover(prs, sd)
        elif vtype == "quote-card": render_quote_card(prs, sd)
        elif vtype == "diagram": render_diagram(prs, sd)
        elif vtype == "screenshot":
            render_screenshot(prs, sd)
            f = SCREENSHOT_MAP.get(sd["page_num"], "")
            if (SCREENSHOTS_DIR / f).exists(): img_ok += 1
            else: img_miss += 1; print(f"  ⚠ 缺失: {f}")
        else: render_text_only(prs, sd)

    intermission(prs)

    for sd in slides_p2:
        vtype = get_visual_type(sd["page_num"])
        stats[vtype] += 1
        if vtype == "section-cover": render_section_cover(prs, sd)
        elif vtype == "quote-card": render_quote_card(prs, sd)
        elif vtype == "diagram": render_diagram(prs, sd)
        elif vtype == "screenshot":
            render_screenshot(prs, sd)
            f = SCREENSHOT_MAP.get(sd["page_num"], "")
            if (SCREENSHOTS_DIR / f).exists(): img_ok += 1
            else: img_miss += 1; print(f"  ⚠ 缺失: {f}")
        else: render_text_only(prs, sd)

    prs.save(str(OUTPUT_PPTX))
    total = len(prs.slides)
    sz = OUTPUT_PPTX.stat().st_size / 1024 / 1024

    print(f"\n✅ 生成完成!")
    print(f"   总页数: {total}")
    print(f"   文件: {sz:.1f} MB")
    print(f"   截图: {img_ok} 找到 / {img_miss} 缺失")
    print(f"   类型分布: {stats}")
    print(f"   输出: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
